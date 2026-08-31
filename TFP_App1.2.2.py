"""
================================================================================
เว็บแอป TFP Executive Summary (Streamlit + Gemini API)
================================================================================
วิธีรัน:
    streamlit run app.py

ต้องมีไฟล์ .streamlit/secrets.toml อยู่ในโฟลเดอร์เดียวกับ app.py นี้ โดยข้างในมี:
    GEMINI_API_KEY = "ใส่ key จริงตรงนี้"

ติดตั้ง library ที่ต้องใช้ (ถ้ายังไม่มี):
    pip install streamlit google-genai pandas numpy statsmodels openpyxl scipy reportlab python-pptx python-docx matplotlib

หมายเหตุ (ส.ค. 2026): โมเดล gemini-2.5-flash เดิมจะถูก retire 16 ต.ค. 2026
ไฟล์นี้อัปเดตให้ใช้ gemini-3.7-flash แล้ว — ก่อนรัน ให้อัปเกรด SDK เป็นเวอร์ชัน
ล่าสุดด้วย: pip install -U google-genai
================================================================================
"""

import streamlit as st
import pandas as pd
import numpy as np
import altair as alt
import base64
import hmac
import os
import re
import math
import time
import warnings
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from io import BytesIO
from datetime import datetime, timezone, timedelta

# เซิร์ฟเวอร์ของแอปมักตั้งเวลาไว้เป็น UTC (ไม่ใช่เวลาไทย) ทำให้เวลาที่แสดงในแอป
# (เช่น "ดึงข้อมูลล่าสุดเมื่อ") ช้ากว่าเวลาจริง 7 ชั่วโมง — ใช้ TH_TZ แทน datetime.now()
# ทุกจุดที่ต้องการแสดงเวลาให้ผู้ใช้เห็น เพื่อให้ตรงกับเวลาประเทศไทย (UTC+7)
TH_TZ = timezone(timedelta(hours=7))


def now_th() -> datetime:
    """คืนค่าเวลาปัจจุบันตามเวลาประเทศไทย (UTC+7) แบบ naive datetime
    (ตัด tzinfo ออกเพื่อให้ยังใช้ .strftime() ต่อได้เหมือนโค้ดเดิมทุกจุด)"""
    return datetime.now(TH_TZ).replace(tzinfo=None)
from google import genai
from google.genai import types
from statsmodels.tsa.arima.model import ARIMA
from statsmodels.tsa.stattools import adfuller
from statsmodels.tools.sm_exceptions import ConvergenceWarning

# ปิด warning จาก statsmodels ตอนลองหาค่า (p,d,q) หลายชุดของ ARIMA (auto-search)
# เพราะบางชุดค่าไม่ converge/ไม่เหมาะกับข้อมูล ซึ่งเป็นเรื่องปกติของการลอง grid search
warnings.simplefilter("ignore", ConvergenceWarning)
warnings.simplefilter("ignore", UserWarning)
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image as RLImage,
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_RIGHT
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from docx import Document
from docx.shared import Pt as DocxPt, Mm as DocxMm, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn as docx_qn
from docx.oxml import OxmlElement as docx_OxmlElement

from TFP import (
    build_model_frame, run_long_run, run_short_run,
    build_coefficient_tables, build_tfpi_yoy_summary, summary_adj_r2,
    adf_report, run_diagnostics, LONG_RUN_VARS, SHORT_RUN_SPEC, DEP_VAR,
)
from data_loader import load_data_gsheet

st.set_page_config(page_title="ระบบวิเคราะห์ผลิตภาพปัจจัยการผลิตรวม", layout="wide")

# ------------------------------------------------------------------------------
# ฟอนต์ — เปลี่ยนหน้าเว็บให้ใช้ 'Prompt' (ฟอนต์ Thai sans-serif ที่เว็บ สอวช./NXPO
# ใช้จริง ตามภาพตัวอย่างที่ผู้ใช้ส่งมา) เป็นฟอนต์หลัก โดยยังคง Sarabun ไว้เป็น
# ฟอนต์สำรอง (ใช้กรณีโหลด Google Fonts ไม่ได้ และยังใช้กับไฟล์ PDF/PPTX ที่ export
# อยู่เดิม เพราะไฟล์ .ttf ของ Sarabun ฝังอยู่ในเครื่องอยู่แล้ว)
#
# หากต้องการเปลี่ยนฟอนต์อีกในอนาคต แก้แค่ตัวแปร FONT_FAMILY ตัวเดียวด้านล่างนี้
# (ต้องเป็นชื่อฟอนต์ที่มีบน Google Fonts และรองรับภาษาไทย เช่น Prompt, Kanit,
# IBM Plex Sans Thai, Noto Sans Thai)
# ------------------------------------------------------------------------------
APP_DIR = os.path.dirname(os.path.abspath(__file__)) if "__file__" in dir() else "."
FONT_DIR = os.path.join(APP_DIR, "fonts")
FONT_FAMILY = "Prompt"  # <-- เปลี่ยนชื่อฟอนต์ตรงนี้ที่เดียวถ้าอยากเปลี่ยนฟอนต์ทั้งเว็บ


def _font_b64(filename: str) -> str:
    with open(os.path.join(FONT_DIR, filename), "rb") as f:
        return base64.b64encode(f.read()).decode()


st.markdown(
    f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family={FONT_FAMILY.replace(" ", "+")}:wght@300;400;500;600;700;800&display=swap');

    /* ใช้ * แทน selector รายชื่อ class เดิม เพื่อให้ครอบคลุมทุก element ของ
       Streamlit จริงๆ (ปุ่ม, input, dropdown, expander, dataframe, tab,
       sidebar ฯลฯ ที่ปกติมี class เฉพาะของตัวเอง override ฟอนต์ default ทับอยู่) */
    html, body, * {{
        font-family: '{FONT_FAMILY}', 'Sarabun', sans-serif !important;
    }}

    /* ไอคอน/ฟอนต์ตัวเลขบางตัวของ Streamlit (เช่น Material Icons/Symbols ใน
       ปุ่ม expander ▶, multiselect X, ฯลฯ) ใช้ font เฉพาะของมันเอง (ligature
       font ที่แปลงคำว่า "keyboard_arrow_right" ให้กลายเป็นรูปลูกศร) ถ้าโดน
       !important ด้านบนบังคับเป็น Prompt จะไม่ใช่ ligature อีกต่อไป กลายเป็น
       ข้อความ "arrow_right" ปนกับ label ตรงๆ จึงต้องกันไว้ให้ครอบคลุมทุกแบบ
       ที่ Streamlit ใช้ (ชื่อ class/attribute เปลี่ยนไปตามเวอร์ชัน) */
    [class*="material-icons" i],
    [class*="material-symbols" i],
    [class*="MaterialIcon" i],
    [data-testid*="Icon" i],
    [data-testid="stExpanderToggleIcon"],
    [data-testid="stIconMaterial"] {{
        font-family: 'Material Symbols Rounded', 'Material Icons', sans-serif !important;
    }}
    </style>
    """,
    unsafe_allow_html=True,
)

try:
    _regular_b64 = _font_b64("Sarabun-Regular.ttf")
    _bold_b64 = _font_b64("Sarabun-Bold.ttf")
    st.markdown(f"""
    <style>
    @font-face {{
        font-family: 'Sarabun';
        src: url(data:font/ttf;base64,{_regular_b64}) format('truetype');
        font-weight: 400;
    }}
    @font-face {{
        font-family: 'Sarabun';
        src: url(data:font/ttf;base64,{_bold_b64}) format('truetype');
        font-weight: 700;
    }}
    </style>
    """, unsafe_allow_html=True)
except FileNotFoundError:
    # ไม่เจอไฟล์ .ttf ของ Sarabun ในเครื่อง -> ไม่เป็นไร เพราะฟอนต์หลักคือ
    # Prompt จาก Google Fonts CDN ด้านบนอยู่แล้ว (Sarabun ใช้เป็นแค่ fallback)
    pass
# ------------------------------------------------------------------------------
# ธีมสี / การ์ด / badge / sidebar ของแดชบอร์ด และแก้ปัญหา multiselect ตัดชื่อ
# ตัวแปรด้วย "..." (ค่าเริ่มต้นของ Streamlit/BaseWeb จำกัดความกว้างของแท็กที่เลือกไว้
# ทำให้ชื่อเต็มของตัวแปรถูกตัดจนอ่านไม่รู้เรื่อง — CSS ด้านล่างแก้ปัญหานี้ไปพร้อมกัน)
# ------------------------------------------------------------------------------
st.markdown("""
<style>
:root {
    --brand-orange: #F2811D;
    --brand-orange-dark: #D96D0F;
    --brand-navy: #16324A;
    --brand-navy-soft: #5B6B7C;
    --bg-page: #F7F5F1;
    --card-border: #E7E1D6;
    --green: #16A34A;
    --amber: #F59E0B;
    --red: #EF4444;
    --blue: #2F6FED;
    --shadow-soft: 0 4px 16px rgba(22,50,74,0.07), 0 1.5px 4px rgba(22,50,74,0.05);
    --shadow-lift: 0 16px 36px rgba(22,50,74,0.13), 0 4px 10px rgba(22,50,74,0.07);
}
.stApp {
    /* พื้นหลังไล่เฉดครีมเบา ๆ + จุดไล่สีส้มจางมากที่มุมบน แทนสีทึบเรียบเดียว
       ให้ความลึกเล็กน้อยโดยไม่ให้สีเพี้ยนไปจากธีม เพื่อให้การ์ดสีขาวด้านบน
       "ลอยเด่น" ขึ้นมาจากพื้นแทนที่จะกลืนไปกับมัน */
    background:
        radial-gradient(1100px 480px at 12% -6%, rgba(242,129,29,0.07), transparent 60%),
        radial-gradient(900px 420px at 100% 0%, rgba(22,50,74,0.045), transparent 55%),
        linear-gradient(180deg, #FAF8F4 0%, var(--bg-page) 320px);
}

/* ----- sidebar: พื้นขาวตามปกติ ไฮไลต์ส้มเฉพาะเมนูที่กำลังเลือกอยู่ ----- */
section[data-testid="stSidebar"] {
    background: #FFFFFF;
    border-right: 1px solid var(--card-border);
}
section[data-testid="stSidebar"] .block-container { padding-top: 1.2rem; }
section[data-testid="stSidebar"] [data-testid="stAlert"] * { color: inherit !important; }

/* ----- การ์ดโลโก้ด้านบนแถบเมนู ----- */
.sidebar-logo-card {
    display: flex; align-items: center; justify-content: center; gap: 24px;
    margin-bottom: 18px;
}

/* ----- ป้ายข้อมูลผู้จัดทำ + โลโก้มหาวิทยาลัย/ภาควิชา + เวอร์ชันแอป —
   วางไว้ท้ายแถบเมนูด้านซ้าย (เล็ก ๆ ไม่เกะกะ ไม่ลอยทับเนื้อหา) ----- */
.corner-badge {
    display: flex;
    flex-direction: column;
    align-items: flex-start;
    gap: 6px;
    padding: 0;
    margin-top: 18px;
    padding-top: 12px;
    border-top: 1px solid var(--card-border);
}
.corner-badge-logos {
    display: flex; align-items: center; gap: 8px; flex-shrink: 0;
}
.corner-badge-text {
    font-size: 0.62rem; line-height: 1.4; color: var(--brand-navy-soft);
    text-align: left;
    width: 100%;
}
.corner-badge-author {
    font-weight: 700; color: var(--brand-navy); font-size: 0.66rem;
}
.corner-badge-version {
    margin-top: 2px; font-weight: 600; color: var(--brand-orange-dark);
}
.sidebar-section-label {
    color: var(--brand-navy-soft) !important; font-size: 0.78rem; font-weight: 700;
    letter-spacing: 0.04em; margin: 4px 0 10px 6px; text-transform: uppercase;
    display: flex; align-items: center; gap: 7px;
}

/* ----- sidebar nav (ปุ่มเมนู หน้าหลัก / Dashboard) -----
   ปกติพื้นขาว ตัวหนังสือสีเข้ม — พอกด (เมนูนั้นกลายเป็นหน้าที่เลือกอยู่)
   พื้นจะเปลี่ยนเป็นสีส้มของแบรนด์ ตัวหนังสือเป็นสีขาว */
section[data-testid="stSidebar"] div[data-testid="stButton"] button {
    justify-content: flex-start !important;
    border-radius: 10px !important;
    font-size: 0.96rem !important;
    padding: 10px 14px !important;
    margin-bottom: 4px;
}
section[data-testid="stSidebar"] div[data-testid="stButton"] button[kind="secondary"] {
    background: #FFFFFF !important;
    border: 1px solid var(--card-border) !important;
    color: var(--brand-navy) !important;
    font-weight: 500 !important;
}
section[data-testid="stSidebar"] div[data-testid="stButton"] button[kind="secondary"]:hover {
    background: #FFF6EC !important;
    border-color: var(--brand-orange) !important;
    color: var(--brand-orange-dark) !important;
}
section[data-testid="stSidebar"] div[data-testid="stButton"] button[kind="primary"] {
    background: var(--brand-orange) !important;
    border: 1px solid var(--brand-orange) !important;
    color: #FFFFFF !important;
    font-weight: 700 !important;
}
section[data-testid="stSidebar"] div[data-testid="stButton"] button[kind="primary"]:hover {
    background: var(--brand-orange-dark) !important;
    border-color: var(--brand-orange-dark) !important;
    color: #FFFFFF !important;
}
/* ตัวหนังสือในปุ่มจริงๆ อยู่ใน <p>/<span> ซ้อนอยู่ข้างใน ต้องกำหนดสีตรงนี้ด้วย
   ไม่งั้นสีที่ตั้งไว้ที่ตัว <button> จะไม่ถูกนำไปใช้ (ปัญหาเดิมที่เจอ) */
section[data-testid="stSidebar"] div[data-testid="stButton"] button[kind="secondary"] p,
section[data-testid="stSidebar"] div[data-testid="stButton"] button[kind="secondary"] span {
    color: var(--brand-navy) !important;
}
section[data-testid="stSidebar"] div[data-testid="stButton"] button[kind="primary"] p,
section[data-testid="stSidebar"] div[data-testid="stButton"] button[kind="primary"] span {
    color: #FFFFFF !important;
}

/* ----- top header ----- */
.app-header {
    display: flex; justify-content: space-between; align-items: center;
    flex-wrap: wrap; gap: 14px; margin-bottom: 18px;
}
.app-header h1 { font-size: 1.65rem; margin: 0; color: var(--brand-navy); font-weight: 800; letter-spacing: -0.01em; overflow-wrap: break-word; }
.app-header p { margin: 4px 0 0 0; color: var(--brand-navy-soft); font-size: 0.92rem; line-height: 1.6; overflow-wrap: break-word; max-width: 68ch; }
.app-header p.app-header-desc { max-width: none; white-space: nowrap; }
.header-chip {
    display: inline-flex; align-items: center; gap: 8px; max-width: 100%;
    background: #FFFFFF; border: 1px solid var(--card-border);
    padding: 8px 14px; border-radius: 12px; font-size: 0.85rem; color: var(--brand-navy-soft);
    box-shadow: 0 2px 6px rgba(15,23,42,0.06), inset 0 1px 0 rgba(255,255,255,0.9);
    transition: box-shadow .15s ease, border-color .15s ease, transform .15s ease;
}
.header-chip:hover { box-shadow: 0 6px 16px rgba(15,23,42,0.1); border-color: #DCD3C2; transform: translateY(-1px); }
.header-chip span { overflow-wrap: break-word; }
.header-chip svg { color: var(--brand-orange-dark); flex-shrink: 0; }

/* ----- เอฟเฟกต์เข้าฉากแบบนุ่ม ๆ (ลูกเล่นเล็ก ๆ ตอนการ์ดปรากฏ ไม่รบกวนสายตา) ----- */
@keyframes tfp-rise {
    from { opacity: 0; transform: translateY(8px); }
    to   { opacity: 1; transform: translateY(0); }
}

/* ----- metric cards ----- */
.metric-card {
    background: linear-gradient(180deg, #FFFFFF 0%, #FFFDFA 100%);
    border: 1px solid var(--card-border); border-radius: 18px;
    padding: 18px 20px; display: flex; align-items: center; gap: 15px; height: 100%;
    box-shadow: var(--shadow-soft), inset 0 1px 0 rgba(255,255,255,0.9);
    transition: box-shadow .18s ease, transform .18s ease, border-color .18s ease;
    animation: tfp-rise .4s ease both;
    position: relative; overflow: hidden;
}
.metric-card:hover { box-shadow: var(--shadow-lift), inset 0 1px 0 rgba(255,255,255,0.9); transform: translateY(-3px); border-color: #E3D8C4; }
.metric-icon {
    width: 46px; height: 46px; border-radius: 50%;
    display: flex; align-items: center; justify-content: center;
    font-size: 1.25rem; flex-shrink: 0; color: #fff;
    background-image: linear-gradient(155deg, rgba(255,255,255,0.28), rgba(255,255,255,0));
    box-shadow: inset 0 0 0 1px rgba(255,255,255,0.22), 0 6px 14px rgba(22,50,74,0.18), 0 0 0 4px rgba(255,255,255,0.6);
    transition: transform .25s ease;
}
.metric-card:hover .metric-icon { transform: scale(1.06) rotate(-4deg); }
.metric-value { font-size: 1.5rem; font-weight: 800; color: var(--brand-navy); line-height: 1.1; letter-spacing: -0.01em; }
.metric-label {
    font-size: 0.8rem; color: var(--brand-navy-soft); margin-top: 2px;
    overflow-wrap: break-word; line-height: 1.45;
}

/* ----- section card ----- */
.section-card {
    background: linear-gradient(180deg, #FFFFFF 0%, #FFFDFA 100%);
    border: 1px solid var(--card-border); border-radius: 18px;
    padding: 11px 22px; margin-bottom: 22px;
    box-shadow: var(--shadow-soft), inset 0 1px 0 rgba(255,255,255,0.9);
    position: relative; overflow: hidden;
    transition: box-shadow .2s ease;
    animation: tfp-rise .45s ease both;
}
.section-card:hover { box-shadow: var(--shadow-lift), inset 0 1px 0 rgba(255,255,255,0.9); }
/* เส้นไล่สีบาง ๆ ด้านบนการ์ด — จุดสังเกตเล็ก ๆ ให้ดูมีมิติขึ้น ไม่แย่งความสนใจจากเนื้อหา */
.section-card::before {
    content: ""; position: absolute; top: 0; left: 0; right: 0; height: 4px;
    background-image: linear-gradient(90deg, var(--brand-orange) 0%, var(--brand-orange-dark) 35%, transparent 100%);
    opacity: 0.9;
}
.section-title { display: flex; align-items: center; gap: 14px; margin-bottom: 0; }
.section-num {
    width: 44px; height: 44px; border-radius: 50%;
    background-image: linear-gradient(155deg, var(--brand-orange), var(--brand-orange-dark));
    color: #fff; font-weight: 800; display: flex; align-items: center; justify-content: center;
    flex-shrink: 0; font-size: 1.2rem;
    box-shadow: inset 0 0 0 1px rgba(255,255,255,0.25), 0 5px 12px rgba(217,109,15,0.35), 0 0 0 4px rgba(255,255,255,0.6);
}
.section-title-text { min-width: 0; }
.section-title h3 {
    margin: 0; color: var(--brand-navy); font-size: 1.22rem; font-weight: 800; letter-spacing: -0.01em;
    line-height: 1.4; overflow-wrap: break-word; word-break: normal;
}

/* ----- badge pill ----- */
.badge-pill {
    display: inline-flex; align-items: center; gap: 7px; white-space: nowrap;
    padding: 4px 13px; border-radius: 999px; font-size: 0.8rem; font-weight: 700;
    border: 1px solid transparent;
}
.badge-pill::before {
    content: ""; width: 6px; height: 6px; border-radius: 50%; flex-shrink: 0;
    background: currentColor;
}
.badge-pass { background: #EAF8EF; color: #158A41; border-color: #CFEEDA; }
.badge-watch { background: #FDF3E1; color: #B9770E; border-color: #F5E1B8; }
.badge-fail { background: #FCEBEA; color: #C0392B; border-color: #F5CFCB; }
.badge-fail::before { animation: tfp-pulse 1.8s ease-in-out infinite; }
@keyframes tfp-pulse {
    0%, 100% { box-shadow: 0 0 0 0 rgba(192,57,43,0.35); }
    50% { box-shadow: 0 0 0 4px rgba(192,57,43,0); }
}

/* ----- ตาราง HTML สำหรับ Diagnostics ----- */
.tfp-table {
    width: 100%; border-collapse: separate; border-spacing: 0; font-size: 0.86rem;
    border-radius: 12px; overflow: hidden; border: 1px solid var(--card-border);
    background: #FFFFFF; box-shadow: var(--shadow-soft);
}
.tfp-table th {
    background-image: linear-gradient(155deg, var(--brand-navy), #0E2436);
    color: #fff; text-align: center; padding: 10px 10px;
    font-weight: 600; letter-spacing: 0.01em; border-right: 1px solid rgba(255,255,255,0.12);
}
.tfp-table th:last-child { border-right: none; }
.tfp-table td { padding: 9px 10px; border-bottom: 1px solid var(--card-border); color: var(--brand-navy-soft); transition: background .12s ease; text-align: center; }
.tfp-table tr:last-child td { border-bottom: none; }
.tfp-table tr:nth-child(odd) td { background: #FFFFFF; }
.tfp-table tr:nth-child(even) td { background: #EEF2F6; }
.tfp-table tr:hover td { background: #FFF6EC; }

/* ----- ตาราง HTML ธีมครีม-ส้ม สำหรับตัวเลขพยากรณ์ ARIMA ----- */
.tfp-table-cream {
    width: 100%; border-collapse: separate; border-spacing: 0; font-size: 0.86rem;
    border-radius: 14px; overflow: hidden; border: 1px solid #F0DCC0;
    background: #FFFDF9; box-shadow: 0 6px 20px rgba(217,109,15,0.12), 0 1.5px 4px rgba(217,109,15,0.08);
}
.tfp-table-cream th {
    background: var(--brand-orange); color: #fff; text-align: center; padding: 11px 12px;
    font-weight: 700; letter-spacing: 0.01em; border-right: 1px solid rgba(255,255,255,0.25);
}
.tfp-table-cream th:last-child { border-right: none; }
.tfp-table-cream td {
    padding: 10px 12px; border-bottom: 1px solid #F3E7D6; color: var(--brand-navy-soft);
    text-align: center; font-variant-numeric: tabular-nums; transition: background .12s ease;
}
.tfp-table-cream td:first-child { font-weight: 700; color: var(--brand-navy); }
.tfp-table-cream tr:last-child td { border-bottom: none; }
.tfp-table-cream tr:nth-child(odd) td { background: #FFFDF9; }
.tfp-table-cream tr:nth-child(even) td { background: #FFF6E9; }
.tfp-table-cream tr:hover td { background: #FFEBD1; }

/* ----- แบนเนอร์ CTA สร้างสรุป AI ----- */
.ai-banner {
    background-image: linear-gradient(135deg, var(--brand-orange) 0%, var(--brand-orange-dark) 100%);
    border-radius: 14px; padding: 14px 20px; color: #fff; margin-bottom: 18px;
    box-shadow: 0 8px 20px rgba(217,109,15,0.22);
    position: relative; overflow: hidden;
}
/* วงกลมจาง ๆ ตกแต่งมุมขวา — ลูกเล่นเล็ก ๆ ให้แบนเนอร์ดูมีมิติ ไม่แย่งตัวหนังสือ */
.ai-banner::after {
    content: ""; position: absolute; right: -30px; top: -40px; width: 140px; height: 140px;
    border-radius: 50%; background: rgba(255,255,255,0.08); pointer-events: none;
}
.ai-banner h3 { margin: 0; font-size: 1.08rem; letter-spacing: -0.01em; overflow-wrap: break-word; line-height: 1.3; }
.ai-banner p { margin: 2px 0 0 0; font-size: 0.85rem; opacity: 0.92; line-height: 1.4; overflow-wrap: break-word; white-space: nowrap; }

@media (max-width: 700px) {
    .ai-banner p { white-space: normal; }
}

/* ----- แก้ปัญหา multiselect ตัดชื่อตัวแปรด้วย "..." ----- */
[data-baseweb="tag"] {
    max-width: none !important; height: auto !important; min-height: 28px;
    white-space: normal !important; background-color: var(--brand-orange) !important;
}
[data-baseweb="tag"] span {
    max-width: none !important; white-space: normal !important;
    overflow: visible !important; text-overflow: clip !important; word-break: break-word !important;
}
div[data-baseweb="select"] > div { flex-wrap: wrap !important; height: auto !important; }
div[data-testid="stMultiSelect"] [data-baseweb="select"] { height: auto !important; }

div[data-testid="stFileUploader"] section { border-radius: 12px; border: 1.5px dashed #D9C2A6; }
button[kind="primary"] { background: var(--brand-orange) !important; border-color: var(--brand-orange) !important; }

/* ----- ปุ่มดาวน์โหลด (st.download_button) ทั้งหมดในแอป -----
   ปกติปุ่มดาวน์โหลดของ Streamlit จะเป็นสไตล์ "secondary" (พื้นขาว ขอบเทาบาง
   ตัวหนังสือเทา) ทำให้ผู้ใช้มองไม่ออกว่ากดได้/เป็นปุ่มดาวน์โหลด — ด้านล่างนี้
   ปรับเป็นพื้นขาว ขอบส้มของแบรนด์ ตัวหนังสือเทาเข้ม เห็นชัดว่าเป็นปุ่มกดได้ */
div[data-testid="stDownloadButton"] button {
    background: #FFFFFF !important;
    border: 2px solid var(--brand-orange) !important;
    color: #374151 !important;
    font-weight: 700 !important;
    border-radius: 10px !important;
    padding: 10px 20px !important;
    box-shadow: var(--shadow-soft);
    transition: all 0.15s ease;
}
div[data-testid="stDownloadButton"] button:hover {
    background: #FFF6EC !important;
    border-color: var(--brand-orange-dark) !important;
    color: #374151 !important;
    transform: translateY(-1px);
    box-shadow: var(--shadow-lift);
}
div[data-testid="stDownloadButton"] button:active {
    transform: translateY(0px);
}
/* ตัวหนังสือในปุ่มจริงๆ อยู่ใน <p>/<span> ซ้อนอยู่ข้างใน ต้องกำหนดสีตรงนี้ด้วย
   ไม่งั้นสีที่ตั้งไว้ที่ตัว <button> จะไม่ถูกนำไปใช้ */
div[data-testid="stDownloadButton"] button p,
div[data-testid="stDownloadButton"] button span {
    color: #374151 !important;
    font-weight: 700 !important;
}
/* เพิ่มคำอธิบายเล็กๆ ใต้ตัวหนังสือหลักของปุ่ม เพื่อบอกชัดเจนว่ากดเพื่อดาวน์โหลด
   (ใช้ ::after ใส่ไว้ที่ตัว <p> ของปุ่ม ไม่ต้องแก้ label ทีละจุดในโค้ด Python) */
div[data-testid="stDownloadButton"] button p::after {
    content: "กดเพื่อดาวน์โหลด";
    display: block;
    font-size: 0.7rem;
    font-weight: 500;
    color: var(--brand-orange-dark);
    margin-top: 2px;
}

/* ----- กันตัวอักษรไทยล้น/ฉีกกลางคำในทุกข้อความของแอป (คำอธิบาย, caption, ตัวเลข) -----
   ปัญหาเดิม: ข้อความไทยยาว ๆ ในกล่องแคบบางจุดล้นกรอบหรือถูกตัดขวางกลางคำ
   วิธีแก้: อนุญาตให้ตัดคำเมื่อจำเป็นเท่านั้น (ไม่บังคับตัดกลางคำถ้ายังพอมีที่บรรทัดปกติ)
   และเพิ่มระยะห่างบรรทัดให้อ่านง่ายขึ้น ลดความรู้สึก "แน่น/รก" ของกล่องข้อความยาว ๆ */
[data-testid="stMarkdownContainer"] p,
[data-testid="stCaptionContainer"],
[data-testid="stCaptionContainer"] p,
.stAlert p,
div[data-testid="stExpander"] p {
    overflow-wrap: break-word;
    line-height: 1.6;
}
[data-testid="stCaptionContainer"] { line-height: 1.55 !important; }

/* จังหวะเข้าฉากของแบนเนอร์ AI ให้เข้าชุดกับการ์ดอื่น ๆ */
.ai-banner { animation: tfp-rise .4s ease both; }

/* ----- theme widget พื้นฐานของ Streamlit ให้เข้าโทนส้ม/กรมท่าของแบรนด์
   (ปกติ slider/checkbox/radio ใช้สีแดงเริ่มต้นของ Streamlit ซึ่งหลุดโทน) ----- */
:root, .stApp { --primary-color: var(--brand-orange); }

/* slider: หัวจับและช่วงที่ลากผ่านเป็นสีส้มแบรนด์แทนสีแดงเริ่มต้น */
div[data-testid="stSlider"] div[role="slider"] {
    background-color: var(--brand-orange) !important;
    border-color: var(--brand-orange) !important;
    box-shadow: 0 0 0 4px rgba(242,129,29,0.15) !important;
}
div[data-testid="stSlider"] div[data-baseweb="slider"] > div > div:nth-child(2) {
    background: var(--brand-orange) !important;
}
div[data-testid="stTickBar"] { display: none; }
div[data-testid="stSlider"] label p { color: var(--brand-navy); font-weight: 600; }

/* selectbox / multiselect: กรอบมนสอดคล้องกับการ์ด และไฮไลต์ส้มตอนโฟกัส */
div[data-baseweb="select"] > div {
    border-radius: 10px !important; border-color: var(--card-border) !important;
    transition: border-color .15s ease, box-shadow .15s ease;
}
div[data-baseweb="select"]:focus-within > div {
    border-color: var(--brand-orange) !important;
    box-shadow: 0 0 0 1px var(--brand-orange) !important;
}
[data-baseweb="menu"] li[aria-selected="true"] { background: #FFF6EC !important; color: var(--brand-orange-dark) !important; }

/* number/text input: โฟกัสสีส้มแทนสีแดงเริ่มต้น */
div[data-testid="stTextInput"] input:focus,
div[data-testid="stNumberInput"] input:focus {
    border-color: var(--brand-orange) !important; box-shadow: 0 0 0 1px var(--brand-orange) !important;
}

/* checkbox / radio: จุด/ติ๊กสีส้มแบรนด์ */
div[data-testid="stCheckbox"] label span[data-checked="true"],
div[data-testid="stRadio"] label span[aria-checked="true"] {
    background-color: var(--brand-orange) !important; border-color: var(--brand-orange) !important;
}

/* container ที่มีเส้นขอบ (st.container(border=True)) ให้โค้งมนเข้าชุดการ์ด */
div[data-testid="stVerticalBlockBorderWrapper"] > div[data-testid="stVerticalBlock"] {
    border-radius: 14px !important;
}

/* expander: โค้งมน มี hover เบา ๆ ให้รู้สึกกดได้ */
div[data-testid="stExpander"] summary {
    border-radius: 12px !important; font-weight: 600; color: var(--brand-navy);
    transition: background .15s ease;
}
div[data-testid="stExpander"] summary:hover { background: #FFF6EC !important; }

/* progress bar / spinner: สีส้มแบรนด์ */
div[data-testid="stSpinner"] > div { border-top-color: var(--brand-orange) !important; }
.stProgress > div > div > div { background-color: var(--brand-orange) !important; }

/* กล่องแจ้งเตือน (info/warning/success/error): โค้งมน มีเงาบาง ๆ เข้าชุดกับการ์ดอื่น ๆ */
div[data-testid="stAlert"] {
    border-radius: 14px !important;
    box-shadow: var(--shadow-soft);
    animation: tfp-rise .35s ease both;
}
</style>
""", unsafe_allow_html=True)

import base64


def img_to_base64(path):
    with open(path, "rb") as f:
        return base64.b64encode(f.read()).decode()


logo1_path = os.path.join(APP_DIR, "สอวช_Logo.png")
logo2_path = os.path.join(APP_DIR, "สวค_Logo.png")
_LOGO1_SIZE = 110  # px — ขนาดโลโก้ สอวช (ปรับแยกต่างหาก)
_LOGO2_SIZE = 76   # px — ขนาดโลโก้ สวค (ปรับแยกต่างหาก)

# โลโก้สถาบันการศึกษา (มหาวิทยาลัย + ภาควิชา) — วางไฟล์ทั้งสองไว้ในโฟลเดอร์
# เดียวกับ app.py นี้ โดยตั้งชื่อไฟล์ตามด้านล่าง (หรือแก้ path ให้ตรงกับไฟล์จริง)
logo3_path = os.path.join(APP_DIR, "kmutnb_logo.png")   # โลโก้ มจพ. (สี่เหลี่ยม)
logo4_path = os.path.join(APP_DIR, "dept_logo.png")     # โลโก้ภาควิชาสถิติประยุกต์ (แนวนอน)
_LOGO3_SIZE = 76    # px — ขนาดโลโก้มหาวิทยาลัย (สี่เหลี่ยมจัตุรัส เท่ากับโลโก้สวค)
_LOGO4_HEIGHT = 46  # px — ความสูงโลโก้ภาควิชา (ตัวนี้เป็นภาพแนวนอน จึงกำหนดแค่ความสูง)
_LOGO4_WIDTH = 180  # px — ความกว้างสูงสุดของกล่องโลโก้ภาควิชา


def _logo_box(size: int) -> str:
    return f'height:{size}px;width:{size}px;display:flex;align-items:center;justify-content:center;'


def _logo_box_wh(width: int, height: int) -> str:
    return f'height:{height}px;width:{width}px;display:flex;align-items:center;justify-content:center;'


logo1_html = (
    f'<div style="{_logo_box(_LOGO1_SIZE)}"><img src="data:image/png;base64,{img_to_base64(logo1_path)}" '
    f'style="max-height:{_LOGO1_SIZE}px;max-width:{_LOGO1_SIZE}px;width:auto;height:auto;object-fit:contain;"></div>'
) if os.path.exists(logo1_path) else ""
logo2_html = (
    f'<div style="{_logo_box(_LOGO2_SIZE)}"><img src="data:image/png;base64,{img_to_base64(logo2_path)}" '
    f'style="max-height:{_LOGO2_SIZE}px;max-width:{_LOGO2_SIZE}px;width:auto;height:auto;object-fit:contain;"></div>'
) if os.path.exists(logo2_path) else ""
logo3_html = (
    f'<div style="{_logo_box(_LOGO3_SIZE)}"><img src="data:image/png;base64,{img_to_base64(logo3_path)}" '
    f'style="max-height:{_LOGO3_SIZE}px;max-width:{_LOGO3_SIZE}px;width:auto;height:auto;object-fit:contain;"></div>'
) if os.path.exists(logo3_path) else ""
logo4_html = (
    f'<div style="{_logo_box_wh(_LOGO4_WIDTH, _LOGO4_HEIGHT)}"><img src="data:image/png;base64,{img_to_base64(logo4_path)}" '
    f'style="max-height:{_LOGO4_HEIGHT}px;max-width:{_LOGO4_WIDTH}px;width:auto;height:auto;object-fit:contain;"></div>'
) if os.path.exists(logo4_path) else ""

# ------------------------------------------------------------------------------
# เวอร์ชันแอป + ข้อมูลผู้จัดทำ — ย้ายมาไว้ตรงนี้ (แทนที่จะอยู่ท้ายไฟล์) เพราะต้อง
# ใช้ประกอบกล่องมุมขวาบน (corner badge) ที่ประกาศไว้ถัดไปด้านล่าง
#
# แก้ข้อมูลผู้จัดทำ/อาจารย์ที่ปรึกษา/ช่องทางติดต่อได้ที่ตัวแปรด้านล่างนี้ที่เดียว
# ------------------------------------------------------------------------------
APP_VERSION = "1.2.2"

AUTHOR_NAME = "นางสาวปรญา ดอกพิกุล"
AUTHOR_PROGRAM = "สาขาวิชาสถิติประยุกต์สำหรับวิทยาการวิเคราะห์ธุรกิจและอุตสาหกรรม ภาควิชาสถิติประยุกต์"
AUTHOR_DEPT = "ภาควิชาสถิติประยุกต์"  # ใช้แสดงในป้ายมุมซ้าย (ตัดชื่อสาขาวิชายาวๆ ออก เหลือแค่ภาค)
AUTHOR_FACULTY = "คณะวิทยาศาสตร์ประยุกต์"
AUTHOR_UNIVERSITY = "มหาวิทยาลัยเทคโนโลยีพระจอมเกล้าพระนครเหนือ"
AUTHOR_YEAR = "2569"
AUTHOR_ADVISOR = ""   # เช่น "อาจารย์ที่ปรึกษา: ผศ.ดร. ชื่อ นามสกุล" — เว้นว่างไว้ถ้ายังไม่ระบุ
AUTHOR_CONTACT = ""   # เช่น "example@email.com" — เว้นว่างไว้ถ้ายังไม่ต้องการเผยแพร่

_footer_lines = [AUTHOR_NAME]
if AUTHOR_DEPT or AUTHOR_FACULTY:
    _footer_lines.append(" ".join(x for x in (AUTHOR_DEPT, AUTHOR_FACULTY) if x))
if AUTHOR_UNIVERSITY:
    _footer_lines.append(AUTHOR_UNIVERSITY)
if AUTHOR_ADVISOR:
    _footer_lines.append(AUTHOR_ADVISOR)

_footer_meta = f"ปีการศึกษา {AUTHOR_YEAR}" if AUTHOR_YEAR else ""
if AUTHOR_CONTACT:
    _footer_meta = f"{_footer_meta} • ติดต่อ: {AUTHOR_CONTACT}" if _footer_meta else f"ติดต่อ: {AUTHOR_CONTACT}"

# ป้ายข้อมูลผู้จัดทำ (fixed html string) — รวมโลโก้มหาวิทยาลัย + ภาควิชา,
# ข้อมูลผู้จัดทำ และเลขเวอร์ชันแอปไว้ด้วยกัน แสดงเล็ก ๆ ท้ายแถบเมนูด้านซ้าย
# ใช้ขนาดย่อมกว่าเดิม (logo3_html/logo4_html) เพราะป้ายนี้ต้องกะทัดรัด
_CORNER_LOGO3_SIZE = 26
_CORNER_LOGO4_WIDTH, _CORNER_LOGO4_HEIGHT = 70, 18
_corner_logo3_html = (
    f'<div style="{_logo_box(_CORNER_LOGO3_SIZE)}"><img src="data:image/png;base64,{img_to_base64(logo3_path)}" '
    f'style="max-height:{_CORNER_LOGO3_SIZE}px;max-width:{_CORNER_LOGO3_SIZE}px;width:auto;height:auto;object-fit:contain;"></div>'
) if os.path.exists(logo3_path) else ""
_corner_logo4_html = (
    f'<div style="{_logo_box_wh(_CORNER_LOGO4_WIDTH, _CORNER_LOGO4_HEIGHT)}"><img src="data:image/png;base64,{img_to_base64(logo4_path)}" '
    f'style="max-height:{_CORNER_LOGO4_HEIGHT}px;max-width:{_CORNER_LOGO4_WIDTH}px;width:auto;height:auto;object-fit:contain;"></div>'
) if os.path.exists(logo4_path) else ""
_corner_logo_html = "".join(h for h in (_corner_logo3_html, _corner_logo4_html) if h)
_corner_badge_html = (
    f'<div class="corner-badge">'
    + (f'<div class="corner-badge-logos">{_corner_logo_html}</div>' if _corner_logo_html else "")
    + f'<div class="corner-badge-text">'
    f'<div class="corner-badge-author">{_footer_lines[0]}</div>'
    + "".join(f'<div>{line}</div>' for line in _footer_lines[1:] if line)
    + (f'<div>{_footer_meta}</div>' if _footer_meta else "")
    + f'<div class="corner-badge-version">v{APP_VERSION}</div>'
    + '</div></div>'
)

# ------------------------------------------------------------------------------
# ชุดไอคอนเส้น (inline SVG) — ใช้แทนอิโมจิสีสันในจุดตกแต่ง UI เพื่อความเรียบหรู
# และสม่ำเสมอของภาพลักษณ์ (ไม่กระทบข้อความ/เนื้อหาใด ๆ ที่แสดงผลอยู่เดิม)
# หมายเหตุ: ใช้ stroke="currentColor" เพื่อให้สีไอคอนไหลตาม CSS `color` ของ
# กล่องแม่โดยอัตโนมัติ (เช่นในวงกลมพื้นสี ไอคอนจะเป็นสีขาวตาม .metric-icon)
# ------------------------------------------------------------------------------
_ICON_PATHS = {
    "check": '<path d="M4 10.5L8 14.5L16 6" stroke-linecap="round" stroke-linejoin="round"/>',
    "alert": (
        '<path d="M10 3L18 17H2L10 3Z" stroke-linejoin="round"/>'
        '<line x1="10" y1="8.3" x2="10" y2="12" stroke-linecap="round"/>'
        '<circle cx="10" cy="14.4" r="0.9" fill="currentColor" stroke="none"/>'
    ),
    "x": '<path d="M5 5L15 15M15 5L5 15" stroke-linecap="round"/>',
    "trend-up": (
        '<polyline points="3,14 8,9 11.5,12.5 17,6" fill="none" stroke-linecap="round" stroke-linejoin="round"/>'
        '<polyline points="12,6 17,6 17,11" fill="none" stroke-linecap="round" stroke-linejoin="round"/>'
    ),
    "trend-down": (
        '<polyline points="3,6 8,11 11.5,7.5 17,14" fill="none" stroke-linecap="round" stroke-linejoin="round"/>'
        '<polyline points="12,14 17,14 17,9" fill="none" stroke-linecap="round" stroke-linejoin="round"/>'
    ),
    "database": (
        '<ellipse cx="10" cy="5" rx="6.5" ry="2.4"/>'
        '<path d="M3.5 5V15C3.5 16.3 6.4 17.4 10 17.4C13.6 17.4 16.5 16.3 16.5 15V5"/>'
        '<path d="M3.5 10C3.5 11.3 6.4 12.4 10 12.4C13.6 12.4 16.5 11.3 16.5 10"/>'
    ),
    "file": (
        '<path d="M6 2.5H12.5L16 6V17C16 17.3 15.8 17.5 15.5 17.5H6C5.7 17.5 5.5 17.3 5.5 17V3C5.5 2.7 5.7 2.5 6 2.5Z" stroke-linejoin="round"/>'
        '<path d="M12.5 2.5V6H16" stroke-linejoin="round"/>'
        '<line x1="7.5" y1="10" x2="13.5" y2="10" stroke-linecap="round"/>'
        '<line x1="7.5" y1="13" x2="13.5" y2="13" stroke-linecap="round"/>'
    ),
    "users": (
        '<circle cx="7.3" cy="7" r="2.6"/>'
        '<path d="M2.5 16.5C2.5 13.6 4.6 11.7 7.3 11.7C10 11.7 12.1 13.6 12.1 16.5" stroke-linecap="round"/>'
        '<circle cx="13.8" cy="6.5" r="2" opacity="0.6"/>'
        '<path d="M12.8 11.6C15.2 11.9 17 13.7 17 16.5" stroke-linecap="round" opacity="0.6"/>'
    ),
    "search": (
        '<circle cx="8.5" cy="8.5" r="5.5"/>'
        '<line x1="13" y1="13" x2="17.5" y2="17.5" stroke-linecap="round"/>'
    ),
    "bulb": (
        '<path d="M10 2.5C6.7 2.5 4.5 4.9 4.5 7.8C4.5 9.7 5.4 11 6.6 12.1C7.2 12.6 7.5 13.3 7.5 14.1V14.7H12.5V14.1C12.5 13.3 12.8 12.6 13.4 12.1C14.6 11 15.5 9.7 15.5 7.8C15.5 4.9 13.3 2.5 10 2.5Z" stroke-linejoin="round"/>'
        '<line x1="7.7" y1="17" x2="12.3" y2="17" stroke-linecap="round"/>'
        '<line x1="8.3" y1="15.4" x2="8.3" y2="14.7"/>'
        '<line x1="11.7" y1="15.4" x2="11.7" y2="14.7"/>'
    ),
    "bars": (
        '<line x1="4" y1="17" x2="4" y2="10" stroke-linecap="round"/>'
        '<line x1="10" y1="17" x2="10" y2="4" stroke-linecap="round"/>'
        '<line x1="16" y1="17" x2="16" y2="7.5" stroke-linecap="round"/>'
    ),
    "clock": (
        '<circle cx="10" cy="10" r="7.2"/>'
        '<path d="M10 6V10L12.6 12" stroke-linecap="round" stroke-linejoin="round"/>'
    ),
    "info": (
        '<circle cx="10" cy="10" r="7.2"/>'
        '<line x1="10" y1="9" x2="10" y2="14" stroke-linecap="round"/>'
        '<circle cx="10" cy="6.3" r="0.9" fill="currentColor" stroke="none"/>'
    ),
}


def icon(name: str, size: int = 18, stroke_width: float = 1.6) -> str:
    """คืนค่า inline SVG ของไอคอนเส้น (ไม่ใช้สีตายตัว — สืบสีจาก CSS `color`
    ของกล่องแม่ผ่าน currentColor) สำหรับแทรกแทนอิโมจิในจุดตกแต่งของหน้าเว็บ"""
    body = _ICON_PATHS[name]
    return (
        f'<svg width="{size}" height="{size}" viewBox="0 0 20 20" fill="none" '
        f'xmlns="http://www.w3.org/2000/svg" stroke="currentColor" '
        f'stroke-width="{stroke_width}" style="display:inline-block;vertical-align:middle;">'
        f'{body}</svg>'
    )


# ------------------------------------------------------------------------------
# ตั้งค่า Gemini จาก secrets.toml
# ------------------------------------------------------------------------------
GEMINI_MODEL = "gemini-3.7-flash"  # รุ่นล่าสุด (13 ส.ค. 2026); เดิมใช้ "gemini-2.5-flash" ซึ่งจะถูก retire 16 ต.ค. 2026

# โควตาฟรีของ Gemini API ต่ำมาก (RPD อาจแค่ 20 ครั้ง/วันในบาง project/tier) และ
# เมื่อชนโควตา Google จะตอบ error รหัส 429 (RESOURCE_EXHAUSTED) กลับมา — ค่าด้านล่าง
# ควบคุมการ retry อัตโนมัติเวลาเจอ 429 ก่อนที่จะยอมแพ้แล้วโชว์ข้อความแจ้งผู้ใช้
_GEMINI_RETRYABLE_STATUS_CODES = {429, 500, 502, 503, 504}
_GEMINI_MAX_ATTEMPTS = 3
_GEMINI_BASE_DELAY = 2.0  # วินาที (จะเว้น 2s, 4s, 8s ตาม exponential backoff)

try:
    gemini_client = genai.Client(api_key=st.secrets["GEMINI_API_KEY"])
except Exception as e:
    st.error(
        "ไม่พบ GEMINI_API_KEY ใน .streamlit/secrets.toml — เพิ่มบรรทัด "
        '`GEMINI_API_KEY = "ใส่ key จริง"` ในไฟล์นั้นก่อนรัน\n\n'
        f"รายละเอียด error จริง: {type(e).__name__}: {e}"
    )
    st.stop()


SYSTEM_PROMPT = """คุณคือนักเศรษฐศาสตร์ที่ทำหน้าที่จัดทำบทสรุปผู้บริหาร (Executive Summary) จากผล
การวิเคราะห์เชิงปริมาณ (quantitative model output) เสนอต่อผู้บริหารระดับสูงของหน่วยงานภาครัฐ
ที่ไม่มีพื้นฐานทางเศรษฐศาสตร์ (เช่น แพทย์ วิศวกร)

โทนภาษา: ใช้ภาษาไทยที่เป็นทางการ สุภาพ เหมาะกับเอกสารราชการ/รายงานนำเสนอผู้บริหารองค์กรภาครัฐ
(หลีกเลี่ยงภาษาพูดหรือคำที่เป็นกันเองเกินไป) แต่ยังต้องอ่านเข้าใจง่าย ไม่ใช้ศัพท์เทคนิคพร่ำเพรื่อ

หลักการเขียน:
- หลีกเลี่ยงศัพท์เทคนิค (เช่น elasticity, coefficient, cointegration, error correction)
  ถ้าจำเป็นต้องใช้ ให้อธิบายความหมายสั้น ๆ ต่อท้ายทันที
- ใช้ภาษาที่เป็นรูปธรรม เปรียบเทียบกับสิ่งที่คนทั่วไปคุ้นเคย
- เน้นตัวเลขสำคัญและนัยเชิงนโยบาย ไม่ต้องอธิบายวิธีการคำนวณทางสถิติ
- เมื่อกล่าวถึงตัวแปรที่มีนัยสำคัญทางสถิติ (p-value < 0.05 หรือดีกว่า) ให้ระบุสั้น ๆ ในทำนอง
  "ผลการวิเคราะห์แบบจำลองแสดงให้เห็นว่าตัวแปร A และ B มีนัยสำคัญทางสถิติที่ระดับความเชื่อมั่น 95%
  สะท้อนว่าหากประเทศไทยสามารถยกระดับตัวแปรดังกล่าวได้ จะส่งผลเชิงบวกต่อการเติบโตของผลิตภาพ (TFP)"
  โดยกล่าวถึงเชิงเทคนิคแค่พอสังเขป ไม่ต้องอธิบายวิธีทดสอบนัยสำคัญ

โครงสร้างรายงาน (Executive Summary) มี 3 ส่วน:
1. สรุปภาพรวมสถานการณ์เศรษฐกิจปัจจุบัน — บอกทิศทางหลักและปัจจัยขับเคลื่อน 2-3 ข้อ
   โดยอ้างอิงค่าจากตาราง TFP ทั้งระยะสั้นและระยะยาวที่ให้มา (ไม่ต้องพิมพ์ตารางซ้ำ
   เพราะมีตารางแสดงแยกต่างหากให้ผู้อ่านดูอยู่แล้ว) และอธิบายว่าตัวเลขนี้แปลว่าอะไร
   ในภาษาง่าย ๆ
2. เปรียบเทียบกับปีก่อนหรือช่วงก่อนหน้า — ระบุทิศทางการเปลี่ยนแปลงและสาเหตุที่
   เป็นไปได้แบบสั้นกระชับ
3. ข้อเสนอแนะเชิงพยากรณ์จากค่าสัมประสิทธิ์ — แปลค่าสัมประสิทธิ์เป็นข้อความเชิง
   นโยบาย/ธุรกิจ พร้อมข้อเสนอแนะเชิงปฏิบัติ 2-3 ข้อ โดยอ้างอิงตัวแปรที่มีนัยสำคัญทางสถิติ
   เป็นหลัก

ความยาวไม่เกิน 2 หน้า A4 ใช้หัวข้อย่อยชัดเจน ห้ามพิมพ์ตาราง markdown (บรรทัดที่ขึ้นต้น
ด้วย |) หรือคัดลอกตัวเลขจากตารางมาเรียงเป็นตารางซ้ำโดยเด็ดขาด ให้เขียนเป็นความเรียง/
bullet point อ้างอิงตัวเลขในเนื้อหาแทน"""


THAI_MONTHS = ["", "มกราคม", "กุมภาพันธ์", "มีนาคม", "เมษายน", "พฤษภาคม", "มิถุนายน",
               "กรกฎาคม", "สิงหาคม", "กันยายน", "ตุลาคม", "พฤศจิกายน", "ธันวาคม"]


def thai_timestamp() -> str:
    now = now_th()
    return (f"{now.day} {THAI_MONTHS[now.month]} {now.year + 543} "
            f"เวลา {now.strftime('%H:%M')} น.")


def _web_summary_text(summary_text: str) -> str:
    """ตัดข้อความสรุปสำหรับแสดงบน "หน้าเว็บ" เท่านั้น (ไม่กระทบไฟล์ Word/สไลด์ที่ยังใช้
    summary_text เต็มเหมือนเดิม) โดยตัดเฉพาะ "บรรทัด" ที่เป็นคำขึ้นต้นแบบหนังสือราชการ
    ซึ่งเรียกผู้บริหาร (เช่น "เรียน ท่านผู้บริหาร...") ทิ้งไปเท่านั้น — เนื้อหาอธิบายผล
    ก่อนหน้าและถัดจากบรรทัดนั้นยังคงอยู่ครบ ไม่ถูกตัดตามไปด้วย"""
    lines = [
        line for line in summary_text.split("\n")
        if not ("เรียน" in line.strip() and "ผู้บริหาร" in line.strip())
    ]
    return "\n".join(lines).strip()


def generate_summary_gemini(lr_res, sr_res, model_df: pd.DataFrame, dep_ln: str) -> str:
    lr_table, sr_table = build_coefficient_tables(lr_res, sr_res)
    yoy_text = build_tfpi_yoy_summary(model_df, dep_ln)

    user_prompt = f"""ผลการรันโมเดล TFP (Total Factor Productivity) ของประเทศไทย
วิธี: Engle-Granger 2-Step (สมการระยะยาว + Restricted Error Correction Model)

--- สมการระยะยาว (Long-run) ---
Adj. R^2 = {summary_adj_r2(lr_res):.4f}
{lr_table.to_string(index=False)}

--- สมการระยะสั้น (Short-run ECM) ---
Adj. R^2 = {summary_adj_r2(sr_res):.4f}
{sr_table.to_string(index=False)}

--- การเปลี่ยนแปลงของผลิตภาพ (TFPI) ปีล่าสุดเทียบปีก่อนหน้า ---
{yoy_text}

โปรดเขียนบทสรุปผู้บริหารตามโครงสร้าง 3 ส่วนที่กำหนด"""

    last_error = None
    for attempt in range(1, _GEMINI_MAX_ATTEMPTS + 1):
        try:
            response = gemini_client.models.generate_content(
                model=GEMINI_MODEL,
                contents=user_prompt,
                config=types.GenerateContentConfig(system_instruction=SYSTEM_PROMPT),
            )
            return response.text
        except Exception as e:
            # google-genai ห่อ error ของ Gemini API ไว้ใน exception ที่มักมี
            # .code หรือ .status_code เป็นรหัส HTTP (เช่น 429) — ดึงออกมาแบบ
            # กันเหนียวหลายทาง เพราะ SDK บางเวอร์ชันตั้งชื่อ attribute ไม่ตรงกัน
            status_code = getattr(e, "code", None) or getattr(e, "status_code", None)
            last_error = e
            if status_code not in _GEMINI_RETRYABLE_STATUS_CODES or attempt == _GEMINI_MAX_ATTEMPTS:
                break
            time.sleep(_GEMINI_BASE_DELAY * (2 ** (attempt - 1)))

    # retry ครบจำนวนครั้งแล้วยังไม่สำเร็จ (หรือเจอ error ที่ retry ซ้ำไปก็ไม่หาย
    # เช่น API key ผิด) — แจ้งผู้ใช้ด้วยข้อความที่เข้าใจง่าย แทนที่จะโยน
    # exception ดิบๆ ออกไปให้แอป crash
    status_code = getattr(last_error, "code", None) or getattr(last_error, "status_code", None)
    if status_code == 429:
        st.error(
            "ขณะนี้มีการเรียกใช้งาน Gemini API ถี่เกินโควตาที่กำหนดไว้ชั่วคราว "
            "(rate limit / quota เต็ม) กรุณารอสักครู่แล้วลองใหม่อีกครั้ง "
            "หากเกิดขึ้นบ่อย ควรพิจารณาเปิดใช้งาน billing เพื่อขยายโควตา "
            "ใน Google AI Studio"
        )
    else:
        st.error(
            "ไม่สามารถสร้างบทสรุปผู้บริหารจาก Gemini API ได้ "
            f"รายละเอียด error: {type(last_error).__name__}: {last_error}"
        )
    st.stop()


# ------------------------------------------------------------------------------
# สร้างไฟล์ PDF (โลโก้ สอวช. + สวค. บนหัวกระดาษ, ฟอนต์ Sarabun, ตาราง TFP)
# ------------------------------------------------------------------------------
_THAI_FONTS_REGISTERED = False


def _register_thai_fonts():
    global _THAI_FONTS_REGISTERED
    if _THAI_FONTS_REGISTERED:
        return
    pdfmetrics.registerFont(TTFont("Sarabun", os.path.join(FONT_DIR, "Sarabun-Regular.ttf")))
    pdfmetrics.registerFont(TTFont("Sarabun-Bold", os.path.join(FONT_DIR, "Sarabun-Bold.ttf")))
    _THAI_FONTS_REGISTERED = True


def _xml_escape(text: str) -> str:
    """แปลงอักขระ &, <, > ในข้อความดิบให้เป็น XML entity ที่ปลอดภัยก่อนแทรกเข้า
    ReportLab Paragraph (ซึ่ง parse เนื้อหาแบบ XML/HTML ย่อย) — ต้องเรียกก่อนแทรก
    markup ของเราเอง (<b>, &nbsp; ฯลฯ) เสมอ

    นี่คือสาเหตุของบั๊ก "R&D" กลายเป็น "R&D;" ที่เจอ: ตัว "&" ดิบที่ไม่ได้ escape
    ทำให้ ReportLab พยายาม parse มันเป็นจุดเริ่มต้นของ XML entity แล้วไล่กิน
    ตัวอักษรถัดไปเรื่อยๆ จนกว่าจะเจอ ';' ตัวแรกที่พบ (ซึ่งอาจมาจาก '&nbsp;' ที่เรา
    แทรกไว้ทีหลังในข้อความเดียวกัน) ทำให้ข้อความระหว่างทางหายไปและเหลือ ';' ค้าง"""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _md_line_to_flowable(line: str, styles):
    """แปลงข้อความ markdown แบบง่าย ๆ ที่ Gemini ส่งกลับมา (##, -, **bold**)
    เป็น Paragraph ของ reportlab; แถวตาราง markdown (|...|) จะถูกข้ามไป เพราะ
    เราวาดตาราง TFP เองแยกต่างหากด้วยข้อมูลจริงจากโมเดล ไม่ใช้ตารางที่ AI พิมพ์มา"""
    line = line.strip()
    if not line or set(line) <= set("|-: "):
        return None
    if line.startswith("|"):
        return None  # แถวตาราง markdown -> ข้าม (มีตารางจริงแยกอยู่แล้ว)
    line = _xml_escape(line)  # ต้อง escape & < > ก่อนแทรก <b>/&nbsp; ของเราเอง
    text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", line)
    if line.startswith("### "):
        return Paragraph(f"<b>{text.lstrip('# ')}</b>", styles["H3TH"])
    if line.startswith("## "):
        return Paragraph(f"<b>{text.lstrip('# ')}</b>", styles["H2TH"])
    if line.startswith("# "):
        return Paragraph(f"<b>{text.lstrip('# ')}</b>", styles["H1TH"])
    if line.startswith(("- ", "• ", "* ")):
        return Paragraph(f"•&nbsp;&nbsp;{text.lstrip('-•* ')}", styles["BodyTH"])
    return Paragraph(text, styles["BodyTH"])



# ------------------------------------------------------------------------------
# ตารางตัวแปร ววน. — รหัสตัวแปรดิบจาก TFP.py -> ชื่อเต็มภาษาไทย + รหัสย่อ
# ลำดับนี้คือลำดับการแสดงผลในตาราง (ตามแบบฟอร์มมาตรฐานของรายงาน)
# ------------------------------------------------------------------------------
VARIABLE_ORDER = [
    "const", "FDI_GDP", "FEE_GDP", "ln_HDI", "ln_RDH_GDP", "RDG_GDP", "RDP_GDP",
    "ln_JOUR_GDP", "ln_PCT_GDP", "ln_PATENT_GDP", "ln_TUM_GDP", "INDUS_GDP",
    "TRADE_GDP", "MKTCOM", "ECM",
]

VARIABLE_LABELS = {
    "const": "ค่าคงที่ : c",
    "FDI_GDP": "สัดส่วนการลงทุนโดยตรงจากต่างประเทศต่อ GDP : FDI/GDP",
    "FEE_GDP": "ค่าธรรมเนียมในการใช้ทรัพย์สินทางปัญญาต่อ GDP : FEE/GDP",
    "ln_HDI": "ดัชนีการพัฒนามนุษย์ : ln(HDI)",
    "ln_RDH_GDP": "สัดส่วนบุคลากรด้าน R&D ต่อประชากรล้านคน : ln(RDH/GDP)",
    "RDG_GDP": "สัดส่วนการลงทุนด้านวิจัยและพัฒนาของภาครัฐต่อ GDP : RDG/GDP",
    "RDP_GDP": "สัดส่วนการลงทุนด้านวิจัยและพัฒนาของภาคเอกชนต่อ GDP : RDP/GDP",
    "ln_JOUR_GDP": "สัดส่วนจำนวนสิ่งพิมพ์ทางวิทยาศาสตร์และเทคนิคต่อ GDP : ln(JOUR/GDP)",
    "ln_PCT_GDP": "จำนวนสนธิสัญญาความร่วมมือด้านสิทธิบัตรต่อ GDP : ln(PCT/GDP)",
    "ln_PATENT_GDP": "สัดส่วนจำนวนสิทธิบัตรต่อ GDP : ln(PATENT/GDP)",
    "ln_TUM_GDP": "สัดส่วนจำนวนอนุสิทธิบัตรต่อ GDP : ln(TUM/GDP)",
    "INDUS_GDP": "สัดส่วนมูลค่าเพิ่มภาคอุตสาหกรรมต่อ GDP : INDUS/GDP",
    "TRADE_GDP": "อัตราการเปิดกว้างทางการค้า : TRADE/GDP",
    "MKTCOM": "ดัชนีความซับซ้อนทางเศรษฐกิจด้านการค้า : MKTCOM",
    "ECM": "ECM",
}


def var_label_with_abbr(code: str) -> str:
    """แสดงชื่อเต็มของตัวแปร แล้ววงเล็บรหัสย่อไว้ข้างหลัง
    เช่น "FDI_GDP" -> "การลงทุนโดยตรงจากต่างประเทศ (FDI/GDP)"
    ถ้าไม่มีในพจนานุกรม VARIABLE_LABELS จะคืนรหัสเดิม

    รองรับป้ายจากตาราง Stationarity (Short-run) ที่มีสัญลักษณ์ transform นำหน้า
    และ/หรือ lag ต่อท้าย เช่น "ΔFDI_GDP" หรือ "Δ²RDG_GDP (t-2)" — แกะสัญลักษณ์/lag
    ออกก่อน หาโค้ดตัวแปรฐานใน VARIABLE_LABELS แล้วประกอบกลับเป็น "Δ ชื่อเต็ม (ย่อ) (t-n)"
    ถ้าหาโค้ดฐานไม่เจอในพจนานุกรม จะคืนป้ายเดิมแบบไม่แปลง (เหมือนพฤติกรรมเดิม)"""
    m = re.match(r"^(Δ²?)([A-Za-z][A-Za-z0-9_]*)(\s*\(t-\d+\))?$", code)
    if m:
        diff_symbol, base_code, lag_part = m.groups()
        base_label = VARIABLE_LABELS.get(base_code)
        if base_label:
            if " : " in base_label:
                full_name, abbr = base_label.split(" : ", 1)
                base_display = f"{full_name} ({abbr})"
            else:
                base_display = base_label
            return f"{diff_symbol} {base_display}{lag_part or ''}"
        return str(code)

    label = VARIABLE_LABELS.get(code)
    if not label:
        return str(code)
    if " : " in label:
        full_name, abbr = label.split(" : ", 1)
        return f"{full_name} ({abbr})"
    return label


def _var_full_name(code: str) -> str:
    """คืนเฉพาะชื่อเต็มภาษาไทยของตัวแปร (ไม่มีรหัสย่อต่อท้าย) ใช้ในข้อความอธิบาย
    เช่น การ์ดผลกระทบที่อยากพูดถึงตัวแปรด้วยชื่อเต็มล้วน ๆ อ่านลื่นกว่า"""
    label = VARIABLE_LABELS.get(code)
    if not label:
        return str(code)
    if " : " in label:
        return label.split(" : ", 1)[0]
    return label


# แกะรหัสตัวแปรดิบ เช่น "d1_FDI_GDP_lag0" -> (base="FDI_GDP", diff=1, lag=0)
# "ECM_lag1" -> (base="ECM", diff=0, lag=1) | "ln_HDI" -> (base="ln_HDI", diff=0, lag=0)
_VAR_CODE_RE = re.compile(r"^(?:d(?P<diff>[12])_)?(?P<base>.+?)(?:_lag(?P<lag>\d+))?$")


def _parse_variable_code(code: str):
    m = _VAR_CODE_RE.match(str(code).strip())
    if not m:
        return str(code).strip(), 0, 0
    base = m.group("base")
    diff = int(m.group("diff")) if m.group("diff") else 0
    lag = int(m.group("lag")) if m.group("lag") else 0
    return base, diff, lag


def _significance_stars(p_value) -> str:
    try:
        p = float(p_value)
    except (TypeError, ValueError):
        return ""
    if p < 0.01:
        return "***"
    if p < 0.05:
        return "**"
    if p < 0.10:
        return "*"
    return ""


def _format_coefficient_cell(coef_value, p_value, diff: int, lag: int, show_significance: bool) -> str:
    try:
        coef_text = f"{float(coef_value):.4f}"
    except (TypeError, ValueError):
        return "-"
    parts = []
    if lag > 0:
        parts.append(f"(t-{lag})")
    parts.append(coef_text)
    if show_significance:
        star = _significance_stars(p_value)
        if star:
            parts.append(star)
    if diff == 1:
        parts.append("Δ")
    elif diff == 2:
        parts.append("Δ²")
    return " ".join(parts)


def _extract_raw_coefficients(df: pd.DataFrame) -> dict:
    """ดึงค่าสัมประสิทธิ์ดิบ (ตัวเลขจริง ไม่ใช่สตริงจัดรูปแบบแบบใน _merge_coefficient_tables)
    จากตาราง lr_table/sr_table ที่ได้จาก build_coefficient_tables()
    คืนค่าเป็น dict: base_var -> {"coef", "p", "diff", "lag", "raw_code"}
    ถ้าตัวแปรฐานเดียวกันมีหลายแถว (เช่น มีหลาย lag) จะเก็บแถวสุดท้ายไว้ ให้ตรงกับ
    ตัวเลขที่แสดงในตารางสัมประสิทธิ์รวมที่ผู้ใช้เห็นอยู่แล้วในหน้าหลัก"""
    var_col, coef_col = df.columns[0], df.columns[1]
    p_col = df.columns[2] if len(df.columns) > 2 else None
    out = {}
    for _, row in df.iterrows():
        base, diff, lag = _parse_variable_code(row[var_col])
        try:
            coef = float(row[coef_col])
        except (TypeError, ValueError):
            continue
        p_val = None
        if p_col is not None:
            try:
                p_val = float(row[p_col])
            except (TypeError, ValueError):
                p_val = None
        out[base] = {"coef": coef, "p": p_val, "diff": diff, "lag": lag, "raw_code": row[var_col]}
    return out


def _merge_coefficient_tables(lr_table: pd.DataFrame, sr_table: pd.DataFrame) -> pd.DataFrame:
    """แปลงตารางระยะยาว/ระยะสั้นจาก TFP.py (คอลัมน์: รหัสตัวแปรดิบ, สัมประสิทธิ์, p-value)
    เป็นตารางเดียวแบบฟอร์มรายงานทางการ — คอลัมน์ระยะยาวแสดงเฉพาะค่าสัมประสิทธิ์
    (ไม่มีดาว/Δ/lag) ส่วนคอลัมน์ระยะสั้นแสดง (t-N) + ค่า + ดาวนัยสำคัญ + Δ/Δ²

    หมายเหตุ: รูปแบบรหัสตัวแปร (d1_/d2_ prefix, _lagN suffix) และรายชื่อตัวแปรเต็ม
    อิงจากตัวอย่างข้อมูลจริงที่ได้รับมา — ถ้า TFP.py มีตัวแปรอื่นเพิ่มเติมนอกเหนือจาก
    VARIABLE_LABELS ด้านบน ให้เพิ่มรายการในดิกชันนารีนั้น"""

    def build_map(df: pd.DataFrame, show_significance: bool) -> dict:
        var_col, coef_col = df.columns[0], df.columns[1]
        p_col = df.columns[2] if len(df.columns) > 2 else None
        out = {}
        for _, row in df.iterrows():
            base, diff, lag = _parse_variable_code(row[var_col])
            p_val = row[p_col] if p_col is not None else None
            out[base] = _format_coefficient_cell(row[coef_col], p_val, diff, lag, show_significance)
        return out

    lr_map = build_map(lr_table, show_significance=False)
    sr_map = build_map(sr_table, show_significance=True)

    rows = []
    for base in VARIABLE_ORDER:
        label = VARIABLE_LABELS.get(base, base)
        lr_val = lr_map.get(base, "n.a." if base == "ECM" else "-")
        sr_val = sr_map.get(base, "-")
        rows.append([label, lr_val, sr_val])

    # กันตกหล่น: ตัวแปรที่มีในข้อมูลจริงแต่ไม่อยู่ใน VARIABLE_ORDER ด้านบน
    known = set(VARIABLE_ORDER)
    extra_bases = [b for b in list(lr_map) + list(sr_map) if b not in known]
    for base in dict.fromkeys(extra_bases):
        rows.append([VARIABLE_LABELS.get(base, base), lr_map.get(base, "-"), sr_map.get(base, "-")])

    return pd.DataFrame(
        rows,
        columns=["ตัวแปร", "ค่าสัมประสิทธิ์\nสมการระยะยาว", "ค่าสัมประสิทธิ์\nสมการระยะสั้น"],
    )


def _docx_set_thai_font(run, size=11, bold=False, color=None, name="Sarabun"):
    """ตั้งฟอนต์ไทยให้ run ของ python-docx ให้ครบทั้ง ascii/hAnsi/cs (จำเป็นสำหรับ
    ภาษาไทยใน Word เพราะ Word แยก font ของอักษรตะวันตก/complex-script ออกจากกัน
    ถ้าตั้งแค่ run.font.name เฉยๆ ตัวอักษรไทยอาจไม่ถูกวาดด้วยฟอนต์ที่ตั้งไว้)"""
    run.font.name = name
    run.font.size = DocxPt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(docx_qn("w:rFonts"))
    if rFonts is None:
        rFonts = docx_OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(docx_qn("w:ascii"), name)
    rFonts.set(docx_qn("w:hAnsi"), name)
    rFonts.set(docx_qn("w:cs"), name)


def _docx_shade_cell(cell, hex_color: str):
    """ใส่สีพื้นหลังให้เซลล์ตาราง python-docx (ไม่มี API สำเร็จรูปให้ใช้ ต้องแทรก
    element w:shd เข้าไปใน tcPr เอง)"""
    tcPr = cell._tc.get_or_add_tcPr()
    shd = docx_OxmlElement("w:shd")
    shd.set(docx_qn("w:val"), "clear")
    shd.set(docx_qn("w:color"), "auto")
    shd.set(docx_qn("w:fill"), hex_color)
    tcPr.append(shd)


def _docx_add_bottom_border(paragraph, color="0F2B46", size=12):
    """วาดเส้นคั่นบาง ๆ ใต้ย่อหน้า (ใช้แทนเส้นคั่นใต้ชื่อเรื่อง เหมือนในไฟล์ PDF เดิม)"""
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = docx_OxmlElement("w:pBdr")
    bottom = docx_OxmlElement("w:bottom")
    bottom.set(docx_qn("w:val"), "single")
    bottom.set(docx_qn("w:sz"), str(size))
    bottom.set(docx_qn("w:space"), "1")
    bottom.set(docx_qn("w:color"), color)
    pBdr.append(bottom)
    pPr.append(pBdr)


def _docx_add_runs_with_bold(paragraph, text: str, size=11, bold=False, color=None):
    """แตกข้อความที่มี **ตัวหนา** แบบ markdown ง่ายๆ ออกเป็นหลาย run ใน python-docx
    (ไม่เหมือน reportlab ที่ parse <b> ให้เองจาก markup เดียว)"""
    parts = re.split(r"(\*\*.+?\*\*)", text)
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**") and len(part) > 4:
            run = paragraph.add_run(part[2:-2])
            _docx_set_thai_font(run, size=size, bold=True, color=color)
        else:
            run = paragraph.add_run(part)
            _docx_set_thai_font(run, size=size, bold=bold, color=color)


def _docx_add_md_line(doc: "Document", line: str, navy_color):
    """แปลงข้อความ markdown แบบง่าย ๆ ที่ Gemini ส่งกลับมา (##, -, **bold**) เป็น
    ย่อหน้าใน python-docx Document — ตรรกะเดียวกับ _md_line_to_flowable ฝั่ง PDF เดิม
    แถวตาราง markdown (|...|) ถูกข้ามไปเช่นกัน เพราะวาดตาราง TFP จริงแยกไว้แล้ว"""
    s = line.strip()
    if not s or set(s) <= set("|-: "):
        return
    if s.startswith("|"):
        return
    if s.startswith("### "):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = DocxPt(4)
        p.paragraph_format.space_after = DocxPt(3)
        _docx_add_runs_with_bold(p, s[4:], size=12, bold=True, color=navy_color)
        return
    if s.startswith("## "):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = DocxPt(6)
        p.paragraph_format.space_after = DocxPt(4)
        _docx_add_runs_with_bold(p, s[3:], size=13, bold=True, color=navy_color)
        return
    if s.startswith("# "):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = DocxPt(6)
        p.paragraph_format.space_after = DocxPt(4)
        _docx_add_runs_with_bold(p, s[2:], size=15, bold=True, color=navy_color)
        return
    if s.startswith(("- ", "• ", "* ")):
        p = doc.add_paragraph()
        p.paragraph_format.left_indent = DocxMm(5)
        p.paragraph_format.space_after = DocxPt(2)
        bullet_run = p.add_run("•  ")
        _docx_set_thai_font(bullet_run, size=11)
        _docx_add_runs_with_bold(p, s.lstrip("-•* "), size=11)
        return
    p = doc.add_paragraph()
    p.paragraph_format.space_after = DocxPt(2)
    _docx_add_runs_with_bold(p, s, size=11)


_MPL_THAI_FONT_LOADED = False


def _register_thai_font_matplotlib():
    """โหลดฟอนต์ Sarabun (ไฟล์เดียวกับที่ใช้ฝั่ง reportlab) เข้า matplotlib font manager
    เพื่อให้ตัวอักษรไทยบนกราฟที่ render เป็นรูปภาพ (สำหรับฝังในไฟล์ Word) อ่านออกได้
    ปกติ ไม่ใช่สี่เหลี่ยมกล่องขาด (tofu) — ทำครั้งเดียวแล้วจำไว้ ไม่ต้องโหลดซ้ำทุกครั้ง"""
    global _MPL_THAI_FONT_LOADED
    if _MPL_THAI_FONT_LOADED:
        return
    try:
        fm.fontManager.addfont(os.path.join(FONT_DIR, "Sarabun-Regular.ttf"))
        fm.fontManager.addfont(os.path.join(FONT_DIR, "Sarabun-Bold.ttf"))
        plt.rcParams["font.family"] = "Sarabun"
    except FileNotFoundError:
        pass
    _MPL_THAI_FONT_LOADED = True


def _build_tfpi_chart_png(tfpi_series: pd.Series) -> bytes:
    """สร้างกราฟเส้นแนวโน้มดัชนี TFPI รายปีแบบเรียบหรู (เส้นสีกรมท่าของแบรนด์ + จุดวงกลม
    สีส้มขอบขาว + พื้นที่ใต้เส้นไล่สีจาง ๆ + เส้น grid บาง ๆ แนวนอน ไม่มีกรอบขวา/บน)
    คืนค่าเป็นรูป PNG พื้นหลังโปร่งใส ความละเอียดสูง สำหรับฝังในไฟล์ Word — python-docx
    ไม่มีกราฟ native แบบ python-pptx จึงต้อง render เป็นรูปภาพแทน"""
    _register_thai_font_matplotlib()
    navy = "#0F2B46"
    orange = "#F2811D"

    x_labels = [str(y) for y in tfpi_series.index]
    y_vals = tfpi_series.values.astype(float)
    x_pos = list(range(len(x_labels)))

    fig, ax = plt.subplots(figsize=(9.4, 3.6), dpi=200)
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")

    ax.plot(x_pos, y_vals, color=navy, linewidth=2.3, zorder=3, solid_capstyle="round")
    y_floor = min(y_vals) - (max(y_vals) - min(y_vals)) * 0.15 if max(y_vals) != min(y_vals) else min(y_vals) - 1
    ax.fill_between(x_pos, y_vals, y_floor, color=navy, alpha=0.08, zorder=1)
    ax.scatter(x_pos, y_vals, color=orange, s=46, zorder=4, edgecolors="white", linewidths=1.1)

    ax.set_ylim(bottom=y_floor)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C7CFD9")
    ax.spines["bottom"].set_color("#C7CFD9")
    ax.grid(axis="y", color="#E9ECF1", linestyle="--", linewidth=0.7, zorder=0)
    ax.tick_params(colors="#5B6B7C", labelsize=9, length=0)

    # โชว์ label แกน x เป็นช่วง ๆ กันแน่นเกินไปเวลามีหลายปี (เหมือนกราฟ Altair ฝั่งเว็บ)
    step = max(1, round(len(x_labels) / 12))
    tick_idx = list(range(0, len(x_labels), step))
    if tick_idx[-1] != len(x_labels) - 1:
        tick_idx.append(len(x_labels) - 1)
    ax.set_xticks(tick_idx)
    ax.set_xticklabels([x_labels[i] for i in tick_idx])

    for spine in ("left", "bottom"):
        ax.spines[spine].set_linewidth(0.8)

    fig.tight_layout()
    buffer = BytesIO()
    fig.savefig(buffer, format="png", transparent=True, bbox_inches="tight")
    plt.close(fig)
    buffer.seek(0)
    return buffer.getvalue()


def build_word_report(summary_text: str, lr_table: pd.DataFrame, sr_table: pd.DataFrame,
                       lr_adj_r2: float, sr_adj_r2: float, model_df: pd.DataFrame) -> bytes:
    """สร้างไฟล์ Word (.docx) ที่แก้ไขได้ แต่จัดหน้าตาให้เหมือนไฟล์ PDF เดิมที่เคย export
    (หัวกระดาษโลโก้คู่ สอวช./สวค. อยู่กึ่งกลาง, ชื่อเรื่อง, ตารางค่าสัมประสิทธิ์, เนื้อหาสรุป
    จาก AI, ท้ายกระดาษ) — แทนที่ build_pdf_report เดิมที่ export เป็น PDF อ่านอย่างเดียว"""
    NAVY = DocxRGBColor(0x0F, 0x2B, 0x46)
    GREY = DocxRGBColor(0x80, 0x80, 0x80)
    WHITE = DocxRGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_ROW = "EEF2F6"
    HEADER_FILL = "0F2B46"

    doc = Document()

    # --- ขนาดกระดาษ A4 + ระยะขอบ ให้เหมือนไฟล์ PDF เดิม (14mm บน/ล่าง, 18mm ซ้าย/ขวา) ---
    section = doc.sections[0]
    section.page_width = DocxMm(210)
    section.page_height = DocxMm(297)
    section.top_margin = DocxMm(14)
    section.bottom_margin = DocxMm(14)
    section.left_margin = DocxMm(18)
    section.right_margin = DocxMm(18)
    usable_width_mm = 210 - 18 - 18

    # ตั้งฟอนต์เริ่มต้นของเอกสารเป็น Sarabun เพื่อให้ข้อความที่ยังไม่ได้ตั้งฟอนต์เอง
    # (เช่น ที่ Word อาจเติมเอง) ไม่หลุดไปเป็นฟอนต์ default อื่น
    normal_style = doc.styles["Normal"]
    normal_style.font.name = "Sarabun"
    normal_style.font.size = DocxPt(11)
    normal_rPr = normal_style.element.get_or_add_rPr()
    normal_rFonts = normal_rPr.find(docx_qn("w:rFonts"))
    if normal_rFonts is None:
        normal_rFonts = docx_OxmlElement("w:rFonts")
        normal_rPr.append(normal_rFonts)
    normal_rFonts.set(docx_qn("w:ascii"), "Sarabun")
    normal_rFonts.set(docx_qn("w:hAnsi"), "Sarabun")
    normal_rFonts.set(docx_qn("w:cs"), "Sarabun")

    # --- หัวกระดาษ: โลโก้ สอวช. + สวค. คู่กันตรงกลาง แล้วชื่อรายงานอยู่บรรทัดถัดมา ---
    logo1_path = os.path.join(APP_DIR, "สอวช_Logo.png")
    logo2_path = os.path.join(APP_DIR, "สวค_Logo.png")
    logo_specs = []
    if os.path.exists(logo1_path):
        logo_specs.append((logo1_path, 50, 25))
    if os.path.exists(logo2_path):
        logo_specs.append((logo2_path, 20, 20))

    if logo_specs:
        logo_table = doc.add_table(rows=1, cols=len(logo_specs))
        logo_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        for cell, (path, w, h) in zip(logo_table.rows[0].cells, logo_specs):
            cell.width = DocxMm(w + 8)
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run()
            run.add_picture(path, width=DocxMm(w), height=DocxMm(h))
        doc.add_paragraph()

    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _docx_add_runs_with_bold(title_p, "บทสรุปผู้บริหาร", size=18, bold=True, color=NAVY)

    subtitle_p = doc.add_paragraph()
    subtitle_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_p.paragraph_format.space_after = DocxPt(6)
    _docx_add_runs_with_bold(
        subtitle_p, "ผลิตภาพการผลิตรวม (Total Factor Productivity) ของประเทศไทย",
        size=12, color=NAVY,
    )
    _docx_add_bottom_border(subtitle_p, color="0F2B46", size=12)

    doc.add_paragraph()

    # --- กราฟแนวโน้มดัชนี TFPI รายปี (รูปภาพเรียบหรู ฝังไว้ก่อนตารางรายละเอียด) ---
    tfpi_series = model_df[DEP_VAR].dropna().sort_index()
    if len(tfpi_series) >= 2:
        chart_title_p = doc.add_paragraph()
        chart_title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        chart_title_p.paragraph_format.space_after = DocxPt(4)
        _docx_add_runs_with_bold(
            chart_title_p, "แนวโน้มดัชนีผลิตภาพการผลิตรวม (TFPI) รายปี",
            size=13, bold=True, color=NAVY,
        )
        chart_png = _build_tfpi_chart_png(tfpi_series)
        chart_p = doc.add_paragraph()
        chart_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        chart_run = chart_p.add_run()
        chart_run.add_picture(BytesIO(chart_png), width=DocxMm(usable_width_mm))
        doc.add_paragraph()

    # --- ตารางผลการทดสอบปัจจัยด้าน ววน. ที่มีต่อผลิตภาพทางเศรษฐกิจไทย ---
    table_title_p = doc.add_paragraph()
    table_title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    table_title_p.paragraph_format.space_after = DocxPt(4)
    _docx_add_runs_with_bold(
        table_title_p,
        "ตารางผลการทดสอบปัจจัยด้านวิทยาศาสตร์ วิจัย และนวัตกรรม (ววน.) "
        "ที่มีต่อผลิตภาพทางเศรษฐกิจไทย",
        size=13, bold=True, color=NAVY,
    )

    combined_table = _merge_coefficient_tables(lr_table, sr_table)
    n_cols = len(combined_table.columns)
    col_props = [0.46, 0.27, 0.27]

    table = doc.add_table(rows=1, cols=n_cols)
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr_cells = table.rows[0].cells
    for i, col_name in enumerate(combined_table.columns):
        cell = hdr_cells[i]
        cell.width = DocxMm(usable_width_mm * col_props[i])
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        lines = str(col_name).split("\n")
        for j, ln in enumerate(lines):
            if j > 0:
                p.add_run().add_break()
            run = p.add_run(ln)
            _docx_set_thai_font(run, size=9, bold=True, color=WHITE)
        _docx_shade_cell(cell, HEADER_FILL)

    for row_i, row_vals in enumerate(combined_table.values.tolist()):
        cells = table.add_row().cells
        for col_i, val in enumerate(row_vals):
            cell = cells[col_i]
            cell.width = DocxMm(usable_width_mm * col_props[col_i])
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT if col_i == 0 else WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(str(val))
            _docx_set_thai_font(run, size=9, bold=(col_i == 0), color=NAVY if col_i == 0 else None)
            if row_i % 2 == 1:
                _docx_shade_cell(cell, LIGHT_ROW)

    doc.add_paragraph()

    # แถว Adj. R² ต่อท้ายตาราง (ไม่ได้อยู่ในตัวตารางสัมประสิทธิ์ เพราะเป็นค่าประเมิน
    # คุณภาพของแต่ละสมการ ไม่ใช่ค่าสัมประสิทธิ์ของตัวแปร)
    r2_table = doc.add_table(rows=1, cols=n_cols)
    r2_table.style = "Table Grid"
    r2_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    r2_values = ["Adj. R²", f"{lr_adj_r2:.4f}", f"{sr_adj_r2:.4f}"]
    for i in range(n_cols):
        cell = r2_table.rows[0].cells[i]
        cell.width = DocxMm(usable_width_mm * col_props[i])
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT if i == 0 else WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(r2_values[i] if i < len(r2_values) else "")
        _docx_set_thai_font(run, size=9, bold=True)

    caption_p = doc.add_paragraph()
    caption_p.paragraph_format.space_before = DocxPt(4)
    _docx_add_runs_with_bold(
        caption_p,
        'หมายเหตุ: ***, **, * หมายถึง นัยสำคัญทางสถิติที่ระดับความเชื่อมั่น 99%, 95%, 90% '
        'ตามลำดับ | "-" หมายถึงตัวแปรที่ไม่ได้อยู่ในสมการนี้ | Δ และ Δ² หมายถึงผลต่างลำดับที่ 1 '
        'และ 2 ตามลำดับ',
        size=8, color=GREY,
    )

    doc.add_paragraph()

    # --- เนื้อหาสรุปจาก AI ---
    # ข้ามบรรทัดคำลงท้าย/ลายเซ็นที่ Gemini อาจแต่งมาเอง (เช่น "ขอแสดงความนับถือ",
    # "(นักเศรษฐศาสตร์)") เพราะเราใส่คำลงท้ายมาตรฐานเองด้านล่างแทน (เหมือนไฟล์ PDF เดิม)
    _signature_markers = ("ขอแสดงความนับถือ", "นักเศรษฐศาสตร์")
    for line in summary_text.split("\n"):
        if any(marker in line for marker in _signature_markers):
            continue
        _docx_add_md_line(doc, line, NAVY)

    doc.add_paragraph()
    footer_p = doc.add_paragraph()
    _docx_add_runs_with_bold(
        footer_p,
        f"จัดทำโดยระบบปัญญาประดิษฐ์ Google Gemini ({GEMINI_MODEL}) — "
        f"สร้างเมื่อวันที่ {thai_timestamp()}",
        size=8, color=GREY,
    )

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ------------------------------------------------------------------------------
# สร้างไฟล์ PowerPoint สำหรับนำเสนอผู้บริหาร (โลโก้มุมบนขวาทุกสไลด์, สรุปจาก AI
# แบ่งเป็นสไลด์ตามหัวข้อ, ตารางค่าสัมประสิทธิ์)
# ------------------------------------------------------------------------------
NAVY = RGBColor(0x0F, 0x2B, 0x46)
SLATE = RGBColor(0x4A, 0x55, 0x61)
LIGHT_GREY = RGBColor(0x9A, 0xA5, 0xB1)
ROW_TINT = RGBColor(0xEE, 0xF2, 0xF6)


def _parse_summary_sections(text: str):
    """แตกข้อความสรุปจาก Gemini (##, **หัวข้อ**, -/•/*) เป็นรายการ (หัวข้อ, [bullet, ...])
    เพื่อนำไปวางเป็นสไลด์แยกตามหัวข้อ ข้ามแถวตาราง markdown และลายเซ็นท้ายเอกสาร
    (ใช้ตรรกะเดียวกับ _md_line_to_flowable ฝั่ง PDF แต่คืนค่าเป็นข้อความล้วน
    เพราะ python-pptx ไม่ต้องการ XML markup แบบ ReportLab)

    หัวข้อย่อยที่มีเลขนำหน้าเดียวกับหัวข้อหลัก (เช่น "3." แล้วตามด้วย "3.1", "3.2")
    จะถูกรวมเข้าเป็นสไลด์เดียวกัน แทนที่จะแยกสไลด์ต่อหัวข้อย่อย เพื่อไม่ให้เนื้อหาที่
    ควรอยู่หน้าเดียวกันถูกกระจายออกเป็นหลายหน้าโดยไม่จำเป็น"""
    signature_markers = ("ขอแสดงความนับถือ", "นักเศรษฐศาสตร์")
    sections = []
    current_title = None
    current_bullets = []
    current_major = None  # เลขหัวข้อหลักของ current_title เช่น "3" จาก "3. ..." หรือ "3.1 ..."

    def _major_number(title: str):
        m = re.match(r"^(\d+)(?:\.\d+)?\.?\s", title + " ")
        return m.group(1) if m else None

    def _flush():
        if current_title is not None or current_bullets:
            sections.append((current_title or "สรุปภาพรวม", current_bullets))

    for raw_line in text.split("\n"):
        line = raw_line.strip()
        if not line or set(line) <= set("|-: ") or line.startswith("|"):
            continue
        if any(marker in line for marker in signature_markers):
            continue
        heading_match = re.match(r"^#{1,3}\s*(.+)$", line)
        bold_only_match = re.match(r"^\*\*(.+?)\*\*:?$", line)
        if heading_match or bold_only_match:
            new_title = (heading_match.group(1) if heading_match else bold_only_match.group(1)).strip()
            new_major = _major_number(new_title)
            if (current_title is not None and new_major is not None
                    and new_major == current_major):
                # หัวข้อย่อยเลขเดียวกับหัวข้อหลักปัจจุบัน — รวมต่อในสไลด์เดิม
                # แทนการขึ้นสไลด์ใหม่ ใส่เป็นบรรทัดหัวข้อย่อยคั่นก่อนกลุ่ม bullet ถัดไป
                current_bullets.append(f"■ {new_title}")
                continue
            _flush()
            current_title = new_title
            current_major = new_major
            current_bullets = []
            continue
        clean = re.sub(r"\*\*(.+?)\*\*", r"\1", line).lstrip("-•* ").strip()
        if clean:
            current_bullets.append(clean)
    _flush()
    return sections


def _estimate_bullet_lines(bullet_text: str, chars_per_line: int = 68) -> float:
    """ประมาณจำนวนบรรทัดที่ bullet หนึ่งข้อความจะตัดคำ (word-wrap) ภายในกล่องข้อความ
    กว้าง ~11.3 นิ้ว ที่ขนาดฟอนต์ 15-16pt — ใช้ประมาณความสูงที่ต้องใช้จริง เพื่อกันเนื้อหา
    ล้นสไลด์ (คำนวณคร่าว ๆ พอเพียงสำหรับแบ่งหน้า ไม่ต้องแม่นยำระดับพิกเซล)"""
    return max(1.0, math.ceil(len(bullet_text) / chars_per_line))


def _chunk_bullets(bullets, max_lines_per_slide: float = 24.0):
    """แบ่งรายการ bullet ยาว ๆ ออกเป็นหลายสไลด์ ไม่ให้เนื้อหาล้นขอบล่างสไลด์เดียว
    แต่ละ bullet คิดน้ำหนักตามจำนวนบรรทัดที่ประมาณไว้ บวกพื้นที่ระยะห่างระหว่างข้อ
    (ปรับเพดานบรรทัดต่อสไลด์ขึ้นจาก 15 เป็น 24 เพื่อให้เนื้อหาของแต่ละหัวข้อ — รวมทั้ง
    หัวข้อย่อยที่ถูกรวมมาจาก _parse_summary_sections — ยุบอยู่ในสไลด์เดียวได้มากที่สุด
    ลดจำนวนสไลด์ '(ต่อ)' ที่ทำให้เนื้อหากระจายเกินความจำเป็น)"""
    chunks, current, lines_used = [], [], 0.0
    for bullet in bullets:
        weight = _estimate_bullet_lines(bullet) + 0.4  # เผื่อระยะห่าง (space_after) ระหว่างข้อ
        if current and lines_used + weight > max_lines_per_slide:
            chunks.append(current)
            current, lines_used = [], 0.0
        current.append(bullet)
        lines_used += weight
    if current:
        chunks.append(current)
    return chunks


def _add_logos_top_right(slide, prs, logo1_path: str, logo2_path: str):
    margin = Inches(0.3)
    height = Inches(0.55)
    paths = [p for p in (logo1_path, logo2_path) if os.path.exists(p)]
    pics = [slide.shapes.add_picture(p, 0, margin, height=height) for p in paths]
    if not pics:
        return
    gap = Inches(0.15)
    total_width = sum(pic.width for pic in pics) + gap * (len(pics) - 1)
    x = prs.slide_width - margin - total_width
    for pic in pics:
        pic.left = int(x)
        x += pic.width + gap


def _add_footer(slide, prs, footer_text: str):
    box = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.45),
        prs.slide_width - Inches(1.0), Inches(0.35),
    )
    tf = box.text_frame
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = footer_text
    run.font.size = Pt(9)
    run.font.color.rgb = LIGHT_GREY
    run.font.name = "Sarabun"


ORANGE = RGBColor(0xF2, 0x81, 0x1D)


def _add_title_accent(slide, prs, x=Inches(0.8), y=Inches(1.22), width=Inches(0.55)):
    """ขีดเส้นสั้น ๆ สีส้มใต้หัวข้อสไลด์ — เพิ่มจุดเน้นเล็ก ๆ ให้สไลด์ดูมีดีไซน์
    มากกว่าปล่อยว่างระหว่างหัวข้อกับเนื้อหา (มินิมอล ไม่ใช่เส้นเต็มความกว้างสไลด์)"""
    bar = slide.shapes.add_shape(1, x, y, width, Pt(3.5))  # 1 = MSO_SHAPE.RECTANGLE
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    bar.shadow.inherit = False


def _add_page_number(slide, prs, page_no: int):
    """ใส่เลขหน้ามุมล่างขวาแบบเล็ก ๆ สีเทาอ่อน ช่วยให้สไลด์ดูครบองค์ประกอบขึ้น"""
    box = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.0), prs.slide_height - Inches(0.45),
        Inches(0.6), Inches(0.35),
    )
    tf = box.text_frame
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(page_no)
    run.font.size = Pt(9)
    run.font.color.rgb = LIGHT_GREY
    run.font.name = "Sarabun"


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_pptx_report(summary_text: str, lr_table: pd.DataFrame, sr_table: pd.DataFrame,
                       lr_adj_r2: float, sr_adj_r2: float, model_df: pd.DataFrame) -> bytes:
    logo1_path = os.path.join(APP_DIR, "สอวช_Logo.png")
    logo2_path = os.path.join(APP_DIR, "สวค_Logo.png")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    footer_text = (f"บทสรุปผู้บริหาร TFP — จัดทำโดยระบบปัญญาประดิษฐ์ Google Gemini "
                    f"({GEMINI_MODEL}) — สร้างเมื่อวันที่ {thai_timestamp()}")

    # --- สไลด์ 1: หน้าปก ---
    slide = _blank_slide(prs)
    _add_logos_top_right(slide, prs, logo1_path, logo2_path)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "บทสรุปผู้บริหาร"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = "Sarabun"

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.85), Inches(11.7), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "ผลิตภาพการผลิตรวม (Total Factor Productivity) ของประเทศไทย"
    run.font.size = Pt(20)
    run.font.color.rgb = SLATE
    run.font.name = "Sarabun"

    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = thai_timestamp()
    run.font.size = Pt(13)
    run.font.color.rgb = LIGHT_GREY
    run.font.name = "Sarabun"

    # --- สไลด์ตาราง ววน. (ย้ายมาไว้เป็นหน้าที่ 2 ต่อจากหน้าปกทันที) ---
    page_no = 1
    combined_table = _merge_coefficient_tables(lr_table, sr_table)
    slide = _blank_slide(prs)
    _add_logos_top_right(slide, prs, logo1_path, logo2_path)
    head = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
    tf = head.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "ตารางผลการทดสอบปัจจัยด้านวิทยาศาสตร์ วิจัย และนวัตกรรม (ววน.)"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = "Sarabun"
    _add_title_accent(slide, prs, y=Inches(1.05))

    n_rows, n_cols = combined_table.shape[0] + 1, combined_table.shape[1]
    tbl_x, tbl_y = Inches(0.6), Inches(1.4)
    tbl_w, tbl_h = Inches(12.1), Inches(5.6)
    gframe = slide.shapes.add_table(n_rows, n_cols, tbl_x, tbl_y, tbl_w, tbl_h)
    table = gframe.table
    table.columns[0].width = int(tbl_w * 0.46)
    table.columns[1].width = int(tbl_w * 0.27)
    table.columns[2].width = int(tbl_w * 0.27)

    for c, col_name in enumerate(combined_table.columns):
        cell = table.cell(0, c)
        cell.text = str(col_name).replace("\n", " ")
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = "Sarabun"

    for r, row in enumerate(combined_table.values.tolist(), start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r % 2 else ROW_TINT
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(10.5)
            run.font.color.rgb = SLATE
            run.font.name = "Sarabun"
            run.font.bold = (c == 0)
    _add_footer(slide, prs, footer_text)
    _add_page_number(slide, prs, page_no)

    # --- สไลด์ 3: ตัวเลขสรุปคุณภาพแบบจำลอง ---
    page_no += 1
    slide = _blank_slide(prs)
    _add_logos_top_right(slide, prs, logo1_path, logo2_path)
    head = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7))
    p = head.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "คุณภาพของแบบจำลองที่ใช้วิเคราะห์"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = "Sarabun"
    _add_title_accent(slide, prs)

    card_specs = [
        ("สมการระยะยาว (Long-run)", lr_adj_r2),
        ("สมการระยะสั้น (Short-run ECM)", sr_adj_r2),
    ]
    card_w, card_h, gap = Inches(5.5), Inches(2.6), Inches(0.6)
    total_w = card_w * 2 + gap
    start_x = (prs.slide_width - total_w) / 2
    for i, (label, value) in enumerate(card_specs):
        x = start_x + i * (card_w + gap)
        y = Inches(2.2)
        card = slide.shapes.add_shape(1, x, y, card_w, card_h)  # 1 = MSO_SHAPE.RECTANGLE
        card.fill.solid()
        card.fill.fore_color.rgb = ROW_TINT
        card.line.color.rgb = LIGHT_GREY
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = f"{value:.4f}"
        r0.font.size = Pt(44)
        r0.font.bold = True
        r0.font.color.rgb = NAVY
        r0.font.name = "Sarabun"
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = f"Adj. R² — {label}"
        r1.font.size = Pt(15)
        r1.font.color.rgb = SLATE
        r1.font.name = "Sarabun"
    _add_footer(slide, prs, footer_text)
    _add_page_number(slide, prs, page_no)

    # --- เตรียมข้อมูล TFPI ไว้แนบกราฟแนวโน้มในสไลด์หัวข้อแรก (ภาพรวม) แทนที่จะ
    # แยกเป็นสไลด์กราฟเดี่ยวต่างหาก (ย้ายมารวมกับเนื้อหาหัวข้อ "1. สรุปภาพรวม...") ---
    tfpi_series = model_df[DEP_VAR].dropna().sort_index()

    # --- สไลด์ตามหัวข้อที่แตกจากสรุปของ AI (ภาพรวม / เปรียบเทียบ / ข้อเสนอแนะ) ---
    # เนื้อหาแต่ละหัวข้ออาจมี bullet ยาวเกินกว่าจะใส่ในสไลด์เดียว จึงแบ่งหน้าอัตโนมัติ
    # ด้วย _chunk_bullets แล้วต่อชื่อหัวข้อด้วย "(ต่อ)" ในสไลด์ถัดไป กันเนื้อหาล้นขอบล่าง
    # (ตัดบรรทัดขึ้นต้นแบบหนังสือราชการที่เรียกผู้บริหาร เช่น "เรียน ท่านผู้บริหาร..."
    # ออกก่อน — ดึกไม่ต้องการให้กลายเป็นสไลด์แยกในชุดสไลด์แสดงผลกลาง)
    sections = _parse_summary_sections(_web_summary_text(summary_text))
    for section_idx, (title, bullets) in enumerate(sections):
        if not bullets:
            continue
        is_first_section = (section_idx == 0) and len(tfpi_series) >= 2
        # หัวข้อ "เปรียบเทียบ..." ใส่กราฟแท่งย้อนหลังประกอบด้วย เพื่อให้เห็นภาพ
        # การเปลี่ยนแปลงชัดเจนขึ้น ไม่ใช่มีแต่ข้อความอย่างเดียว
        show_compare_chart = (not is_first_section) and ("เปรียบเทียบ" in title) and len(tfpi_series) >= 2
        chunks = _chunk_bullets(bullets)
        for chunk_idx, chunk in enumerate(chunks):
            slide_title = title if chunk_idx == 0 else f"{title} (ต่อ)"
            page_no += 1
            slide = _blank_slide(prs)
            _add_logos_top_right(slide, prs, logo1_path, logo2_path)
            head = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.0), Inches(0.9))
            tf = head.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = slide_title
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = NAVY
            run.font.name = "Sarabun"
            _add_title_accent(slide, prs)

            # เฉพาะสไลด์แรกของหัวข้อภาพรวม (แนบกราฟแนวโน้ม) หรือหัวข้อเปรียบเทียบ
            # (แนบกราฟแท่งย้อนหลัง): ลดความกว้างกล่องข้อความลงเพื่อเปิดพื้นที่ฝั่งขวา
            put_line_chart_here = is_first_section and chunk_idx == 0
            put_bar_chart_here = show_compare_chart and chunk_idx == 0
            put_chart_here = put_line_chart_here or put_bar_chart_here
            body_w = Inches(6.6) if put_chart_here else Inches(11.3)
            body = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), body_w, Inches(5.15))
            tf = body.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(chunk):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(10)
                bullet_run = p.add_run()
                bullet_run.text = "●  "
                bullet_run.font.size = Pt(15)
                bullet_run.font.bold = True
                bullet_run.font.color.rgb = ORANGE
                bullet_run.font.name = "Sarabun"
                text_run = p.add_run()
                text_run.text = bullet
                text_run.font.size = Pt(15)
                text_run.font.color.rgb = SLATE
                text_run.font.name = "Sarabun"

            if put_line_chart_here:
                # กราฟแนวโน้มดัชนี TFPI รายปีเต็มช่วง (เดิมเคยแยกเป็นสไลด์กราฟต่างหาก
                # ย้ายมาแนบในสไลด์หัวข้อแรกแทน)
                chart_data = CategoryChartData()
                chart_data.categories = [str(y) for y in tfpi_series.index]
                chart_data.add_series("TFPI", tuple(float(v) for v in tfpi_series.values))
                gframe = slide.shapes.add_chart(
                    XL_CHART_TYPE.LINE_MARKERS,
                    Inches(7.8), Inches(1.6), Inches(4.6), Inches(5.15),
                    chart_data,
                )
                chart = gframe.chart
                chart.has_legend = False
                chart.has_title = True
                chart.chart_title.text_frame.text = "แนวโน้มดัชนี TFPI รายปี"
                ttl_run = chart.chart_title.text_frame.paragraphs[0].runs[0]
                ttl_run.font.size = Pt(12)
                ttl_run.font.bold = True
                ttl_run.font.color.rgb = NAVY
                ttl_run.font.name = "Sarabun"
                series = chart.plots[0].series[0]
                series.smooth = True
                series.format.line.color.rgb = NAVY
                series.format.line.width = Pt(2.0)
                series.marker.style = XL_MARKER_STYLE.CIRCLE
                series.marker.size = 6
                series.marker.format.fill.solid()
                series.marker.format.fill.fore_color.rgb = RGBColor(0xF2, 0x81, 0x1D)
                series.marker.format.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                series.marker.format.line.width = Pt(1.0)
                chart.category_axis.tick_labels.font.size = Pt(9)
                chart.category_axis.tick_labels.font.name = "Sarabun"
                chart.category_axis.format.line.color.rgb = LIGHT_GREY
                chart.value_axis.tick_labels.font.size = Pt(9)
                chart.value_axis.tick_labels.font.name = "Sarabun"
                chart.value_axis.tick_labels.number_format = "0.00"
                chart.value_axis.tick_labels.number_format_is_linked = False
                chart.value_axis.has_major_gridlines = True
                chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE9, 0xEC, 0xF1)
                chart.value_axis.major_gridlines.format.line.width = Pt(0.75)
                chart.value_axis.format.line.color.rgb = LIGHT_GREY
            elif put_bar_chart_here:
                recent = tfpi_series.tail(6)
                chart_data = CategoryChartData()
                chart_data.categories = [str(y) for y in recent.index]
                chart_data.add_series("TFPI", tuple(float(v) for v in recent.values))
                gframe = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    Inches(7.8), Inches(1.6), Inches(4.6), Inches(5.15),
                    chart_data,
                )
                chart = gframe.chart
                chart.has_legend = False
                chart.has_title = True
                chart.chart_title.text_frame.text = "TFPI ย้อนหลัง (ล่าสุด)"
                ttl_run = chart.chart_title.text_frame.paragraphs[0].runs[0]
                ttl_run.font.size = Pt(12)
                ttl_run.font.bold = True
                ttl_run.font.color.rgb = NAVY
                ttl_run.font.name = "Sarabun"
                series = chart.plots[0].series[0]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = ORANGE
                series.format.line.fill.background()
                chart.category_axis.tick_labels.font.size = Pt(9)
                chart.category_axis.tick_labels.font.name = "Sarabun"
                chart.category_axis.format.line.color.rgb = LIGHT_GREY
                chart.value_axis.tick_labels.font.size = Pt(9)
                chart.value_axis.tick_labels.font.name = "Sarabun"
                chart.value_axis.tick_labels.number_format = "0.00"
                chart.value_axis.tick_labels.number_format_is_linked = False
                chart.value_axis.has_major_gridlines = True
                chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE9, 0xEC, 0xF1)
                chart.value_axis.major_gridlines.format.line.width = Pt(0.75)
                chart.value_axis.format.line.color.rgb = LIGHT_GREY

            _add_footer(slide, prs, footer_text)
            _add_page_number(slide, prs, page_no)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
# ------------------------------------------------------------------------------
# session state: ชุดตัวแปรที่ "ใช้งานจริง" ในสมการของ session นี้ แยกจาก
# LONG_RUN_VARS/SHORT_RUN_SPEC ที่ import มาจาก TFP.py (ค่า default ในไฟล์โค้ด
# จะไม่ถูกแก้ไข) — คณะวิจัยปรับชุดตัวแปรนี้ได้ผ่าน UI ด้านล่าง (ดู expander
# "ปรับตัวแปรในสมการ") โดยทุกการปรับต้องกรอกเหตุผล + พิมพ์คำยืนยันก่อนมีผลจริง
# ------------------------------------------------------------------------------
if "active_long_run_vars" not in st.session_state:
    st.session_state.active_long_run_vars = list(LONG_RUN_VARS)
if "active_short_run_spec" not in st.session_state:
    st.session_state.active_short_run_spec = list(SHORT_RUN_SPEC)
if "var_audit_log" not in st.session_state:
    st.session_state.var_audit_log = []  # แต่ละรายการ: เวลา/ตัดออก/เพิ่มกลับ/เหตุผล

# ------------------------------------------------------------------------------
# แถบด้านข้าง: โลโก้ + เมนูนำทาง + ช่องอัปโหลดข้อมูล
# ------------------------------------------------------------------------------
# หน้าที่มีในเมนู — สลับให้ "Dashboard" ขึ้นเป็นเมนูบนสุด/หน้าเริ่มต้น เพราะเป็น
# หน้าที่เปิดให้บุคคลภายนอกเข้าชมได้ ส่วนหน้าบทสรุปผู้บริหารเดิม (เดิมชื่อ "หน้าหลัก")
# ย้ายลงมาอยู่ล่าง Dashboard และเปลี่ยนชื่อเป็น "สำหรับคณะวิจัยเท่านั้น" พร้อมล็อก
# ด้วย username/password (ดูส่วน RESEARCH_USERNAME/RESEARCH_PASSWORD ด้านล่าง)
NAV_ITEMS = [
    ("Dashboard", "dashboard"),
    ("สำหรับคณะวิจัยเท่านั้น", "home"),
]

if "page" not in st.session_state:
    st.session_state.page = "dashboard"

# บัญชีสำหรับเข้าหน้า "สำหรับคณะวิจัยเท่านั้น" — อ่านจาก .streamlit/secrets.toml
# แทนการฝังในโค้ดตรงๆ (ไฟล์ secrets.toml ต้องมี RESEARCH_USERNAME / RESEARCH_PASSWORD
# และไม่ควร push ขึ้น git — ใส่ไว้ใน .gitignore ด้วย)
try:
    RESEARCH_USERNAME = st.secrets["RESEARCH_USERNAME"]
    RESEARCH_PASSWORD = st.secrets["RESEARCH_PASSWORD"]
    _research_login_config_error = None
except (KeyError, FileNotFoundError):
    RESEARCH_USERNAME = None
    RESEARCH_PASSWORD = None
    _research_login_config_error = (
        "ยังไม่ได้ตั้งค่า RESEARCH_USERNAME / RESEARCH_PASSWORD ใน .streamlit/secrets.toml "
        "จึงยังเข้าหน้านี้ไม่ได้ — กรุณาเพิ่มค่าทั้งสองในไฟล์ secrets.toml ก่อน"
    )
if "research_authenticated" not in st.session_state:
    st.session_state.research_authenticated = False

with st.sidebar:
    _logo_divider_height = max(_LOGO1_SIZE, _LOGO2_SIZE) - 12
    st.markdown(
        f'<div class="sidebar-logo-card">'
        f'{logo1_html}'
        f'<div style="width:1px;height:{_logo_divider_height}px;background:var(--card-border);"></div>'
        f'{logo2_html}</div>',
        unsafe_allow_html=True,
    )
    # หมายเหตุ: โลโก้สถาบันการศึกษา (มหาวิทยาลัย + ภาควิชา) ที่เคยแสดงเป็นแถวที่ 2
    # ตรงนี้ ถูกย้ายไปรวมกับข้อมูลผู้จัดทำและเลขเวอร์ชันแอปในกล่องมุมขวาบนแทนแล้ว
    for label, page_key in NAV_ITEMS:
        is_active = st.session_state.page == page_key
        if st.button(
            label,
            key=f"nav_{page_key}",
            use_container_width=True,
            type="primary" if is_active else "secondary",
        ):
            st.session_state.page = page_key
            st.rerun()

    # ปุ่มออกจากระบบ — โชว์เฉพาะตอนล็อกอินเข้าหน้า "สำหรับคณะวิจัยเท่านั้น" อยู่แล้ว
    if st.session_state.research_authenticated:
        if st.button("ออกจากระบบ (คณะวิจัย)", use_container_width=True, key="nav_logout"):
            st.session_state.research_authenticated = False
            st.session_state.page = "dashboard"
            st.rerun()

    st.markdown("---")
    st.markdown(
        f'<div class="sidebar-section-label">{icon("database", 14, 1.6)}<span>ข้อมูล</span></div>',
        unsafe_allow_html=True,
    )
    if st.button("คลิกดึงข้อมูลอัตโนมัติ", use_container_width=True):
        st.session_state.pop("gsheet_load_error", None)
        try:
            with st.spinner("กำลังดึงข้อมูลอัตโนมัติ..."):
                st.session_state.gsheet_raw_df = load_data_gsheet()
            st.session_state.gsheet_loaded_at = now_th()
        except Exception as e:
            st.session_state.gsheet_load_error = str(e)
            st.session_state.pop("gsheet_raw_df", None)
    if st.session_state.get("gsheet_load_error"):
        st.error(f"ดึงข้อมูลไม่สำเร็จ: {st.session_state.gsheet_load_error}")
    elif "gsheet_raw_df" in st.session_state:
        st.success(f"ดึงข้อมูลล่าสุดเมื่อ {st.session_state.gsheet_loaded_at.strftime('%H:%M:%S')}")
    st.caption("ดึงข้อมูล → รันโมเดล → สรุปผลอัตโนมัติ")

    # ป้ายข้อมูลผู้จัดทำ + โลโก้มหาวิทยาลัย/ภาควิชา + เวอร์ชันแอป — วางไว้ท้าย
    # แถบเมนูด้านซ้าย (เล็ก ๆ ไม่เกะกะ) แทนที่จะลอยทับเนื้อหาแบบเดิม
    st.markdown(_corner_badge_html, unsafe_allow_html=True)

# ------------------------------------------------------------------------------
# รันโมเดล (ถ้ามีไฟล์อัปโหลด) — คำนวณผลลัพธ์ทั้งหมดไว้ก่อน เพื่อนำไปแสดงในการ์ด
# สรุปสถานะที่หัวหน้าเพจ (metric cards) และในแต่ละหมวดด้านล่าง
# ------------------------------------------------------------------------------
result_ready = False
diag_table_display = None
n_pass = n_watch = n_fail = 0
adj_r2_lr = adj_r2_sr = None
vars_customized = False
lr_raw_map = {}
sr_raw_map = {}

if "gsheet_raw_df" in st.session_state:
    active_lr_vars = st.session_state.active_long_run_vars
    active_sr_spec = st.session_state.active_short_run_spec
    default_sr_bases = [c for c, _, _ in SHORT_RUN_SPEC]
    active_sr_bases = [c for c, _, _ in active_sr_spec]

    with st.spinner("กำลังรันโมเดล..."):
        raw = st.session_state.gsheet_raw_df
        model_df = build_model_frame(raw)
        dep_ln = "ln_" + DEP_VAR
        lr_res, lr_resid = run_long_run(model_df, dep_ln, active_lr_vars)
        sr_res = run_short_run(model_df, dep_ln, active_sr_spec, lr_resid)

    if sr_res is None:
        st.error("รันสมการระยะสั้นไม่สำเร็จ (พารามิเตอร์ >= observations) — ตรวจสอบข้อมูลนำเข้า")
    else:
        lr_table, sr_table = build_coefficient_tables(lr_res, sr_res)
        combined_table = _merge_coefficient_tables(lr_table, sr_table)
        lr_raw_map = _extract_raw_coefficients(lr_table)
        sr_raw_map = _extract_raw_coefficients(sr_table)
        adj_r2_lr = summary_adj_r2(lr_res)
        adj_r2_sr = summary_adj_r2(sr_res)
        vars_customized = (active_lr_vars != list(LONG_RUN_VARS) or active_sr_bases != default_sr_bases)

        try:
            diag_table = run_diagnostics(model_df, dep_ln, active_lr_vars, lr_res, sr_res, lr_resid,
                                          short_run_spec=active_sr_spec)
            diag_table_display = diag_table.copy()
            diag_table_display["รายการ"] = diag_table_display["รายการ"].apply(var_label_with_abbr)
            n_fail = int((diag_table["สถานะ"] == "🔴 ไม่ผ่าน").sum())
            n_watch = int(diag_table["สถานะ"].str.startswith("🟡").sum())
            n_pass = int(len(diag_table) - n_fail - n_watch)
        except Exception as e:
            st.info(f"ไม่สามารถรันตารางตรวจสอบข้อสมมติฐานได้ครบทุกรายการ: {e}")

        result_ready = True

# ------------------------------------------------------------------------------
# เนื้อหาของแต่ละหน้า แยกตามเมนูด้านซ้าย: "สำหรับคณะวิจัยเท่านั้น" (บทสรุปผู้บริหาร
# เดิม — ต้องล็อกอินก่อนถึงจะเห็น) กับ "Dashboard" (กราฟแนวโน้ม TFP + ตัวแปรอิสระ
# ที่เปิดให้บุคคลภายนอกเข้าชมได้โดยไม่ต้องล็อกอิน)
# ------------------------------------------------------------------------------
if st.session_state.page == "home":
    if not st.session_state.research_authenticated:
        # หน้าล็อกอิน — แสดงแทนเนื้อหาบทสรุปผู้บริหารจนกว่าจะกรอก user/password ถูกต้อง
        st.markdown(
            f'<div class="section-card" style="max-width:420px;margin:40px auto;">'
            f'<div class="section-title"><div class="section-num">🔒</div>'
            f'<div class="section-title-text"><h3>สำหรับคณะวิจัยเท่านั้น</h3></div></div>'
            f'<p style="color:var(--brand-navy-soft);font-size:0.9rem;margin-top:-6px;">'
            f'กรุณาเข้าสู่ระบบด้วยบัญชีคณะวิจัยเพื่อดูหน้านี้</p></div>',
            unsafe_allow_html=True,
        )
        _login_col = st.columns([1, 1.4, 1])[1]
        with _login_col:
            if _research_login_config_error:
                st.error(_research_login_config_error)
            else:
                with st.form("research_login_form", clear_on_submit=False):
                    login_user = st.text_input("ชื่อผู้ใช้ (Username)")
                    login_pass = st.text_input("รหัสผ่าน (Password)", type="password")
                    login_submitted = st.form_submit_button("เข้าสู่ระบบ", use_container_width=True)
                if login_submitted:
                    # ใช้ hmac.compare_digest แทน == ธรรมดา เพื่อลดความเสี่ยงจาก
                    # timing attack (เดารหัสผ่านจากเวลาที่ใช้เทียบสตริง)
                    user_ok = hmac.compare_digest(
                        login_user.encode("utf-8"), RESEARCH_USERNAME.encode("utf-8")
                    )
                    pass_ok = hmac.compare_digest(
                        login_pass.encode("utf-8"), RESEARCH_PASSWORD.encode("utf-8")
                    )
                    if user_ok and pass_ok:
                        st.session_state.research_authenticated = True
                        st.rerun()
                    else:
                        st.error("ชื่อผู้ใช้หรือรหัสผ่านไม่ถูกต้อง")
        st.stop()

    # ------------------------------------------------------------------------------
    # หัวเพจ (header) — ชื่อรายงาน + สถานะไฟล์ที่อัปโหลด + ป้ายผู้ใช้งาน
    # ------------------------------------------------------------------------------
    header_right = (
        f'<div class="header-chip">{icon("database", 14, 1.5)}<span>TFP-Data — อัปเดตล่าสุด '
        f'{st.session_state.gsheet_loaded_at.strftime("%d/%m/%Y %H:%M")}</span></div>'
        if "gsheet_raw_df" in st.session_state else
        f'<div class="header-chip">{icon("file", 14, 1.5)}<span>ยังไม่ได้ดึงข้อมูล</span></div>'
    )
    st.markdown(
        f"""
        <div class="app-header">
            <div>
                <h1>ระบบวิเคราะห์ผลิตภาพปัจจัยการผลิตรวมมหภาคด้วยแบบจำลองเศรษฐมิติ</h1>
                <p class="app-header-desc">และสร้างรายงานสรุปผลสำหรับผู้บริหารด้วยปัญญาประดิษฐ์เพื่อสนับสนุนการติดตามผลและประเมินผลนโยบายด้านวิทยาศาสตร์ วิจัย และนวัตกรรม (ววน.) ของสอวช.</p>
                <p class="app-header-desc" style="margin-top:4px;"><strong>ขอบเขต:</strong> ระบบนี้วิเคราะห์เฉพาะ<strong>ระดับเศรษฐกิจมหภาค (Macro Level)</strong> ด้วยแบบจำลอง Error-Correction Model (ECM) เท่านั้น</p>
            </div>
        </div>
        <div style="display:flex; justify-content:flex-end; margin-bottom:10px;">
            {header_right}
        </div>
        """,
        unsafe_allow_html=True,
    )


    def metric_card(bg, icon, value, label):
        return (
            f'<div class="metric-card"><div class="metric-icon" style="background:{bg};">{icon}</div>'
            f'<div><div class="metric-value">{value}</div><div class="metric-label">{label}</div></div></div>'
        )


    m1, m2, m3, m4, m5 = st.columns(5)
    with m1:
        st.markdown(metric_card("var(--green)", icon("check", 21, 2), n_pass if result_ready else "-", "ผ่านเกณฑ์"), unsafe_allow_html=True)
    with m2:
        st.markdown(metric_card("var(--amber)", icon("alert", 21, 1.6), n_watch if result_ready else "-", "พิจารณาเพิ่มเติม"), unsafe_allow_html=True)
    with m3:
        st.markdown(metric_card("var(--red)", icon("x", 21, 2), n_fail if result_ready else "-", "ไม่ผ่านเกณฑ์"), unsafe_allow_html=True)
    with m4:
        r2_display = f"{adj_r2_lr:.4f}" if adj_r2_lr is not None else "-"
        st.markdown(metric_card("var(--blue)", icon("trend-up", 21, 1.8), r2_display, "Adj. R² (ระยะยาว)"), unsafe_allow_html=True)
    with m5:
        r2_sr_display = f"{adj_r2_sr:.4f}" if adj_r2_sr is not None else "-"
        st.markdown(metric_card("var(--blue)", icon("trend-down", 21, 1.8), r2_sr_display, "Adj. R² (ระยะสั้น)"), unsafe_allow_html=True)

    st.write("")

    if not result_ready:
        st.info("คลิกเพื่อดึงข้อมูลอัตโนมัติจากแถบด้านซ้ายเพื่อเริ่มต้นการวิเคราะห์")
    else:
        if vars_customized:
            st.info(
                "ℹ️ ผลลัพธ์ด้านล่างนี้รันด้วย **ชุดตัวแปรที่คณะวิจัยปรับไว้** ไม่ใช่ค่า default ในไฟล์โค้ด "
                "— ดูรายละเอียดและเหตุผลได้ที่ประวัติการปรับตัวแปรด้านล่าง"
            )

        # ================= หมวด 1: ผลการทดสอบปัจจัย ววน. =================
        st.markdown(
            '<div class="section-card"><div class="section-title">'
            '<div class="section-num">1</div>'
            '<div class="section-title-text"><h3>ผลการทดสอบปัจจัยด้านวิทยาศาสตร์ วิจัย และนวัตกรรม (ววน.) '
            'ที่มีต่อผลิตภาพทางเศรษฐกิจไทย</h3></div>'
            '</div>',
            unsafe_allow_html=True,
        )
        def _style_coef_cell(value) -> str:
            """ใส่สีให้ตัวเลขสัมประสิทธิ์ตามระดับนัยสำคัญทางสถิติ (ดาว ***/**/*)
            เพื่อให้อ่านตารางง่ายขึ้นด้วยสายตา ไม่ต้องไล่หาดาวทีละช่อง"""
            s = str(value)
            if s in ("-", "nan", "None", ""):
                return '<span style="color:#B9C0CB;">-</span>'
            if "***" in s:
                return f'<span style="color:var(--green);font-weight:700;">{s}</span>'
            if "**" in s:
                return f'<span style="color:var(--brand-orange-dark);font-weight:700;">{s}</span>'
            if "*" in s:
                return f'<span style="color:var(--amber);font-weight:600;">{s}</span>'
            return s

        combined_header_html = "".join(f"<th>{c}</th>" for c in combined_table.columns)
        combined_rows_html = "".join(
            "<tr>" + "".join(
                f'<td style="font-weight:600;color:var(--brand-navy);">{v}</td>'
                if i == 0 else f"<td>{_style_coef_cell(v)}</td>"
                for i, v in enumerate(row)
            ) + "</tr>"
            for row in combined_table.values.tolist()
        )
        st.markdown(
            f'<div style="overflow-x:auto;"><table class="tfp-table"><thead><tr>'
            f'{combined_header_html}</tr></thead><tbody>{combined_rows_html}</tbody></table></div>',
            unsafe_allow_html=True,
        )
        st.caption(
            f"Adj. R² (สมการระยะยาว) = {adj_r2_lr:.4f}  |  Adj. R² (สมการระยะสั้น) = {adj_r2_sr:.4f}"
        )
        st.caption(
            'หมายเหตุ: ***, **, * หมายถึง นัยสำคัญทางสถิติที่ระดับความเชื่อมั่น 99%, 95%, 90% '
            'ตามลำดับ | "-" หมายถึงตัวแปรที่ไม่ได้อยู่ในสมการนี้ | Δ และ Δ² หมายถึงผลต่างลำดับที่ 1 '
            'และ 2 ตามลำดับ'
        )
        st.download_button(
            "⬇️ ดาวน์โหลดตาราง (.csv)",
            data=combined_table.to_csv(index=False).encode("utf-8-sig"),
            file_name="ตาราง_ววน.csv",
            mime="text/csv",
            key="dl_combined_table",
        )
        st.markdown('</div>', unsafe_allow_html=True)

        # ================= หมวด 2: ตรวจสอบข้อสมมติฐาน (Diagnostics) =================
        st.markdown(
            '<div class="section-card"><div class="section-title">'
            '<div class="section-num">2</div>'
            '<div class="section-title-text"><h3>ตรวจสอบข้อสมมติฐานของแบบจำลอง (Diagnostics)</h3>'
            '</div></div>',
            unsafe_allow_html=True,
        )
        st.warning(
            "⚠️ **อ่านก่อนใช้ตารางนี้**: สถานะด้านล่างเป็น *สัญญาณเตือนเบื้องต้นสำหรับผู้ทำวิจัย* "
            "เท่านั้น ไม่ใช่คำตัดสินสุดท้ายทางสถิติ — ห้ามอ่านแค่สี/สถานะแล้วสรุปทันทีโดยไม่เข้าใจ "
            "นัยของการทดสอบแต่ละตัว เกณฑ์ผ่าน/ไม่ผ่าน (เช่น VIF > 10, p < 0.05) เป็นกฎเบื้องต้นทั่วไป "
            "ไม่ใช่กฎตายตัวทางทฤษฎี ควรอ่านคอลัมน์ **หมายเหตุ** ประกอบทุกครั้ง และปรึกษาผู้เชี่ยวชาญ "
            "ด้านเศรษฐมิติก่อนนำไปสรุปในรายงานฉบับจริง"
        )
        if diag_table_display is not None:
            def _status_badge(s):
                if s.startswith("🟢"):
                    return '<span class="badge-pill badge-pass">🟢 ผ่าน</span>'
                if s.startswith("🟡"):
                    return '<span class="badge-pill badge-watch">🟡 พิจารณาเพิ่มเติม</span>'
                return '<span class="badge-pill badge-fail">🔴 ไม่ผ่าน</span>'

            rows_html = "".join(
                "<tr>" + "".join(
                    f"<td>{_status_badge(v) if col == 'สถานะ' else v}</td>"
                    for col, v in zip(diag_table_display.columns, row)
                ) + "</tr>"
                for row in diag_table_display.values.tolist()
            )
            header_html = "".join(f"<th>{c}</th>" for c in diag_table_display.columns)
            st.markdown(
                f'<div style="overflow-x:auto;"><table class="tfp-table"><thead><tr>{header_html}</tr></thead>'
                f'<tbody>{rows_html}</tbody></table></div>',
                unsafe_allow_html=True,
            )
            if n_fail:
                st.caption(f"🔴 มี {n_fail} รายการที่ไม่ผ่านเกณฑ์ทั่วไป และ 🟡 {n_watch} รายการที่ก้ำกึ่ง/ต้องพิจารณาเพิ่มเติม")
            elif n_watch:
                st.caption(f"🟡 มี {n_watch} รายการที่ก้ำกึ่ง/ต้องพิจารณาเพิ่มเติม — ไม่มีรายการที่ไม่ผ่านชัดเจน")
            else:
                st.caption("🟢 ทุกรายการผ่านเกณฑ์ทั่วไปเบื้องต้น")
            st.download_button(
                "⬇️ ดาวน์โหลดตาราง Diagnostics (.csv)",
                data=diag_table_display.to_csv(index=False).encode("utf-8-sig"),
                file_name="diagnostics_TFP.csv",
                mime="text/csv",
                key="dl_diag_table",
            )
        st.markdown('</div>', unsafe_allow_html=True)

        # ================= หมวด 3: ปรับตัวแปรในสมการ (สำหรับงานวิจัย) =================
        st.markdown(
            '<div class="section-card"><div class="section-title">'
            '<div class="section-num">3</div>'
            '<div class="section-title-text"><h3>ปรับตัวแปรในสมการ (สำหรับงานวิจัย)</h3>'
            '</div></div>',
            unsafe_allow_html=True,
        )
        with st.expander("🛠️ ปรับตัวแปรในสมการ (สำหรับคณะวิจัย)", expanded=False):
            st.caption(
                "ใช้ส่วนนี้เมื่อพิจารณาจากตาราง Diagnostics ด้านบนแล้วเห็นว่าควรตัด/เพิ่มตัวแปร "
                "กลับเข้าสมการ (เช่น VIF สูงเกินไป) การปรับที่นี่จะไม่แก้ไขไฟล์ TFP.py — มีผลเฉพาะ "
                "รอบการใช้งานนี้เท่านั้น และทุกครั้งที่ปรับจะถูกบันทึกไว้ในประวัติด้านล่างพร้อมเหตุผล"
            )

            all_lr_vars = list(LONG_RUN_VARS)
            new_lr_vars = st.multiselect(
                "ตัวแปรในสมการระยะยาว (Long-run)",
                options=all_lr_vars,
                default=active_lr_vars,
                format_func=var_label_with_abbr,
                key="ms_lr_vars",
            )

            all_sr_bases = [c for c, _, _ in SHORT_RUN_SPEC]
            new_sr_bases = st.multiselect(
                "ตัวแปรในสมการระยะสั้น (Short-run ECM)",
                options=all_sr_bases,
                default=active_sr_bases,
                format_func=var_label_with_abbr,
                key="ms_sr_vars",
            )

            reason = st.text_area(
                "เหตุผลของการปรับ (จำเป็นต้องกรอกก่อนยืนยัน)",
                placeholder='เช่น "ln_HDI มี VIF=68.5 สูงเกินเกณฑ์ และมีสหสัมพันธ์กับ MKTCOM สูงถึง 0.991 '
                            'คณะวิจัยจึงมีมติให้ตัด ln_HDI ออกจากสมการระยะยาว"',
                key="var_change_reason",
            )

            removed_lr = sorted(set(active_lr_vars) - set(new_lr_vars))
            added_lr = sorted(set(new_lr_vars) - set(active_lr_vars))
            removed_sr = sorted(set(active_sr_bases) - set(new_sr_bases))
            added_sr = sorted(set(new_sr_bases) - set(active_sr_bases))
            has_change = bool(removed_lr or added_lr or removed_sr or added_sr)

            if has_change:
                change_parts = []
                if removed_lr:
                    change_parts.append(f"ตัดออก (ระยะยาว): {', '.join(var_label_with_abbr(v) for v in removed_lr)}")
                if removed_sr:
                    change_parts.append(f"ตัดออก (ระยะสั้น): {', '.join(var_label_with_abbr(v) for v in removed_sr)}")
                if added_lr:
                    change_parts.append(f"เพิ่มกลับ (ระยะยาว): {', '.join(var_label_with_abbr(v) for v in added_lr)}")
                if added_sr:
                    change_parts.append(f"เพิ่มกลับ (ระยะสั้น): {', '.join(var_label_with_abbr(v) for v in added_sr)}")
                st.info("การเปลี่ยนแปลงที่จะเกิดขึ้นถ้ายืนยัน: " + " | ".join(change_parts))

                CONFIRM_PHRASE = "ยืนยันการปรับตัวแปร"
                confirm_text = st.text_input(
                    f'พิมพ์คำว่า "{CONFIRM_PHRASE}" ให้ตรงเป๊ะเพื่อยืนยัน (ป้องกันการกดพลาด)',
                    key="confirm_phrase_input",
                )
                reason_ok = reason.strip() != ""
                confirm_ok = confirm_text.strip() == CONFIRM_PHRASE
                if not reason_ok:
                    st.caption("⚠️ ต้องกรอกเหตุผลก่อนจึงจะยืนยันได้")
                if confirm_text and not confirm_ok:
                    st.caption("⚠️ ข้อความยืนยันไม่ตรงกับที่กำหนด กรุณาพิมพ์ให้ตรงเป๊ะ")

                if st.button("✅ ยืนยันและรันโมเดลใหม่ด้วยตัวแปรชุดนี้",
                              disabled=not (reason_ok and confirm_ok)):
                    st.session_state.var_audit_log.append({
                        "เวลา": thai_timestamp(),
                        "ตัดออก (ระยะยาว)": ", ".join(removed_lr) or "-",
                        "ตัดออก (ระยะสั้น)": ", ".join(removed_sr) or "-",
                        "เพิ่มกลับ (ระยะยาว)": ", ".join(added_lr) or "-",
                        "เพิ่มกลับ (ระยะสั้น)": ", ".join(added_sr) or "-",
                        "เหตุผล": reason.strip(),
                    })
                    st.session_state.active_long_run_vars = new_lr_vars
                    st.session_state.active_short_run_spec = [
                        spec for spec in SHORT_RUN_SPEC if spec[0] in new_sr_bases
                    ]
                    st.success("บันทึกและปรับตัวแปรแล้ว กำลังรันโมเดลใหม่...")
                    st.rerun()
            else:
                st.caption("ยังไม่มีการเปลี่ยนแปลงจากชุดตัวแปรที่ใช้อยู่ในขณะนี้")

            if active_lr_vars != list(LONG_RUN_VARS) or active_sr_bases != default_sr_bases:
                if st.button("↩️ คืนค่าเริ่มต้นทั้งหมด (ตามที่กำหนดในโค้ด TFP.py)"):
                    st.session_state.active_long_run_vars = list(LONG_RUN_VARS)
                    st.session_state.active_short_run_spec = list(SHORT_RUN_SPEC)
                    st.rerun()

            if st.session_state.var_audit_log:
                st.markdown("**ประวัติการปรับตัวแปร (Audit log)**")
                audit_df = pd.DataFrame(st.session_state.var_audit_log)
                st.dataframe(audit_df, use_container_width=True, hide_index=True)
                st.download_button(
                    "📥 ดาวน์โหลดประวัติการปรับตัวแปร (.csv)",
                    data=audit_df.to_csv(index=False).encode("utf-8-sig"),
                    file_name="audit_log_ตัวแปรโมเดล.csv",
                    mime="text/csv",
                )
        st.markdown('</div>', unsafe_allow_html=True)

        # ================= หมวด 4: สร้าง Executive Summary ด้วย AI =================
        st.markdown(
            '<div class="ai-banner"><div><h3> สร้างรายงานสรุปและสไลด์นำเสนอผู้บริหารอัตโนมัติด้วยปัญญาประดิษฐ์ </h3>'
            '<p>สรุปผลการวิเคราะห์ พร้อมข้อเสนอแนะเชิงนโยบายอย่างชัดเจน ดาวน์โหลดได้ทั้ง Word และ PPTX</p></div></div>',
            unsafe_allow_html=True,
        )
        if st.button("📑คลิกเพื่อสร้างรายงานสรุปและสไลด์นำเสนอผู้บริหารอัตโนมัติด้วยปัญญาประดิษฐ์ ", type="primary"):
            with st.spinner("กำลังสรุปผลอัตโนมัติ..."):
                try:
                    summary_text = generate_summary_gemini(lr_res, sr_res, model_df, dep_ln)
                    # กันเหนียว: ตัดบรรทัดตาราง markdown (ขึ้นต้นด้วย |) ที่ Gemini อาจ
                    # แอบใส่มาแม้จะสั่งห้ามแล้วใน SYSTEM_PROMPT เพราะตารางจริงแสดงแยกไว้
                    # ด้านบนแล้ว ไม่ต้องการให้ซ้ำ/โชว์รหัสตัวแปรดิบอีกรอบ (เว้นบรรทัดว่างไว้
                    # ตามเดิม ไม่ตัดออก เพื่อไม่ให้ย่อหน้าติดกัน)
                    def _is_markdown_table_row(line: str) -> bool:
                        s = line.strip()
                        if not s:
                            return False
                        return s.startswith("|") or set(s) <= set("|-: ")

                    summary_text = "\n".join(
                        line for line in summary_text.split("\n")
                        if not _is_markdown_table_row(line)
                    )
                    st.markdown(_web_summary_text(summary_text))
                    footer = (f"*จัดทำโดยระบบปัญญาประดิษฐ์ Google Gemini "
                              f"({GEMINI_MODEL}) — สร้างเมื่อวันที่ {thai_timestamp()}*")
                    st.markdown("---")
                    st.caption(footer)

                    full_output = summary_text + "\n\n---\n" + footer
                    word_bytes = build_word_report(
                        summary_text, lr_table, sr_table,
                        summary_adj_r2(lr_res), summary_adj_r2(sr_res), model_df,
                    )

                    pptx_bytes = build_pptx_report(
                        summary_text, lr_table, sr_table,
                        summary_adj_r2(lr_res), summary_adj_r2(sr_res), model_df,
                    )

                    dl_col1, dl_col2, dl_col3 = st.columns(3)
                    with dl_col1:
                        st.download_button(
                            "📄 ดาวน์โหลดสรุปเป็น Word (มีโลโก้ แก้ไขได้)",
                            data=word_bytes,
                            file_name="สรุปผู้บริหาร_TFP.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        )
                    with dl_col2:
                        st.download_button(
                            "ดาวน์โหลดสไลด์นำเสนอผู้บริหาร (.pptx)",
                            data=pptx_bytes,
                            file_name="สไลด์นำเสนอผู้บริหาร_TFP.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )
                    with dl_col3:
                        st.download_button(
                            "ดาวน์โหลดสรุปเป็น .txt",
                            data=full_output,
                            file_name="executive_summary.txt",
                        )
                except Exception as e:
                    st.error(f"เรียก Gemini ไม่สำเร็จ: {e}")

elif st.session_state.page == "dashboard":
    st.markdown(
        """
        <div class="app-header">
            <div>
                <h1>Dashboard</h1>
                <p>แนวโน้มผลิตภาพการผลิตรวม (TFP) และตัวแปรอิสระ จากข้อมูลที่ดึงมาล่าสุด</p>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    def _nice_line_chart(series: pd.Series, color: str = "#F2811D", height: int = 340):
        """สร้างกราฟเส้นด้วย Altair แทน st.line_chart เดิม เพื่อให้ดูสวยและอ่านง่าย
        ขึ้นกว่าเดิม: เส้นโค้งมน มีพื้นที่ใต้เส้นแบบไล่สีจาง ๆ, เส้น grid บางๆ,
        และ label บนแกนปีที่สุ่มแสดงเป็นช่วง ๆ (ไม่ยัดทุกปีจนอ่านไม่ออก) พร้อม
        tooltip บอกปีและค่าที่ชี้เมื่อเอาเมาส์ไปวาง — และปีบนแกน x เป็น string
        (ordinal) เสมอ กัน Vega-Lite ตีความเป็นตัวเลขแล้วใส่ , คั่นหลักพัน"""
        s = series.copy()
        s.index = s.index.map(str)
        df = s.reset_index()
        df.columns = ["ปี", "ค่า"]

        n = len(df)
        step = max(1, round(n / 12))
        tick_vals = df["ปี"].iloc[::step].tolist()
        if df["ปี"].iloc[-1] not in tick_vals:
            tick_vals.append(df["ปี"].iloc[-1])

        base = alt.Chart(df).encode(
            x=alt.X(
                "ปี:O", sort=None, title=None,
                axis=alt.Axis(values=tick_vals, labelAngle=0, grid=False,
                               domain=False, tickColor="#E9ECF1",
                               labelColor="#5B6B7C", labelFontSize=11, labelPadding=6),
            ),
            y=alt.Y(
                "ค่า:Q", title=None,
                axis=alt.Axis(grid=True, gridColor="#EEF1F5", gridDash=[3, 3],
                               domain=False, tickColor="#E9ECF1",
                               labelColor="#5B6B7C", labelFontSize=11),
            ),
            tooltip=[
                alt.Tooltip("ปี:O", title="ปี"),
                alt.Tooltip("ค่า:Q", title="ค่า", format=".4f"),
            ],
        )
        area = base.mark_area(
            interpolate="monotone",
            line=False,
            color=alt.Gradient(
                gradient="linear",
                stops=[
                    alt.GradientStop(color=color, offset=0),
                    alt.GradientStop(color="#FFFFFF", offset=1),
                ],
                x1=1, x2=1, y1=1, y2=0,
            ),
            opacity=0.35,
        )
        line = base.mark_line(
            interpolate="monotone", color=color, strokeWidth=2.6,
            point=alt.OverlayMarkDef(filled=True, size=32, color=color, stroke="#FFFFFF", strokeWidth=1.6),
        )
        chart = (
            (area + line)
            .properties(height=height, padding={"left": 4, "right": 10, "top": 8, "bottom": 4})
            .configure_view(strokeWidth=0)
            .configure_axis(labelFont=FONT_FAMILY, titleFont=FONT_FAMILY)
        )
        st.altair_chart(chart, use_container_width=True)

    @st.cache_data(show_spinner=False)
    def _auto_arima_forecast(series: pd.Series, periods: int,
                              max_p: int = 4, max_d: int = 2, max_q: int = 4):
        """หาโมเดล ARIMA(p,d,q) ที่เหมาะกับข้อมูลที่สุดด้วยวิธี grid search
        (ลองทุกชุด p,d,q ในช่วงที่กำหนด แล้วเลือกชุดที่ค่า AIC ต่ำที่สุด — AIC ยิ่งต่ำ
        ยิ่งหมายถึงโมเดลอธิบายข้อมูลได้ดีโดยไม่ซับซ้อนเกินจำเป็น) จากนั้นพยากรณ์ล่วงหน้า
        `periods` ปี พร้อมช่วงความเชื่อมั่น 95%

        คืนค่า (forecast_df, order) โดย forecast_df มี index เป็นปีในอนาคต และ
        คอลัมน์ mean / lower / upper ส่วน order คือ (p, d, q) ที่เลือกใช้จริง
        ถ้าหาโมเดลที่ fit ได้ไม่สำเร็จเลย จะ fallback เป็น random walk with drift (0,1,0)"""
        y = series.astype(float).values

        # หาลำดับ differencing (d) ที่เหมาะสมก่อนด้วย ADF test (Augmented Dickey-Fuller)
        # แทนที่จะปล่อยให้ AIC เป็นตัวเลือก d เอง เพราะ AIC เปรียบเทียบข้าม d ต่างกัน
        # ไม่ได้ตรงๆ (ข้อมูลที่มีแนวโน้ม/ไม่ stationary มักได้โมเดล d=0 ที่ AIC ต่ำ
        # หลอกๆ จากการฟิตแบบ ARMA แต่พอพยากรณ์ระยะยาวค่าจะไหลกลับไปหาค่าเฉลี่ยของ
        # อนุกรมทั้งหมดแทนที่จะไปตามแนวโน้มจริง ทำให้ค่าพยากรณ์รูดฮวบผิดปกติ)
        def _select_d(vals, max_diff):
            d = 0
            cur = vals.copy()
            while d < max_diff:
                try:
                    pvalue = adfuller(cur, autolag="AIC")[1]
                except Exception:
                    break
                if pvalue < 0.05:
                    break
                cur = np.diff(cur)
                d += 1
            return d

        fixed_d = _select_d(y, max_d)

        best_aic = np.inf
        best_order = None
        best_fit = None
        for p in range(0, max_p + 1):
            for d in (fixed_d,):
                for q in range(0, max_q + 1):
                    if p == 0 and q == 0:
                        continue
                    try:
                        fit = ARIMA(y, order=(p, d, q)).fit()
                        if np.isfinite(fit.aic) and fit.aic < best_aic:
                            best_aic = fit.aic
                            best_order = (p, d, q)
                            best_fit = fit
                    except Exception:
                        continue

        if best_fit is None:
            # กันเหนียว: ถ้าไม่มีชุด (p,d,q) ไหน fit ได้เลย ใช้ random walk with
            # drift แทน (โมเดลพื้นฐานที่สุด ยังพยากรณ์แนวโน้มต่อได้เสมอ)
            best_fit = ARIMA(y, order=(0, 1, 0), trend="t").fit()
            best_order = (0, 1, 0)

        fc = best_fit.get_forecast(steps=periods)
        summary = fc.summary_frame(alpha=0.05)
        last_year = int(series.index.max())
        future_years = [last_year + i for i in range(1, periods + 1)]
        forecast_df = pd.DataFrame(
            {
                "mean": summary["mean"].values,
                "lower": summary["mean_ci_lower"].values,
                "upper": summary["mean_ci_upper"].values,
            },
            index=future_years,
        )
        return forecast_df, best_order

    def _nice_line_chart_with_forecast(hist_series: pd.Series, forecast_df: pd.DataFrame,
                                        color: str = "#F2811D", forecast_color: str = "#2F6FED",
                                        height: int = 340):
        """เหมือน _nice_line_chart แต่ต่อเส้นพยากรณ์ (เส้นประสีน้ำเงิน) และแถบ
        ช่วงความเชื่อมั่น 95% (พื้นที่สีน้ำเงินจาง ๆ) ต่อจากข้อมูลจริงให้ในกราฟเดียวกัน"""
        hist = hist_series.copy()
        hist.index = hist.index.map(int)
        years_hist = list(hist.index)
        years_fc = list(forecast_df.index)
        all_years = years_hist + years_fc
        year_order = [str(y) for y in all_years]

        df = pd.DataFrame({"ปี": year_order, "ปี_num": all_years})
        df["ข้อมูลจริง"] = df["ปี_num"].map(hist.to_dict())

        # เชื่อมจุดสุดท้ายของข้อมูลจริงเข้ากับเส้นพยากรณ์ ไม่ให้เส้นขาดตอน
        last_year, last_val = years_hist[-1], float(hist.iloc[-1])
        fc_mean = {last_year: last_val, **forecast_df["mean"].to_dict()}
        fc_lower = {last_year: last_val, **forecast_df["lower"].to_dict()}
        fc_upper = {last_year: last_val, **forecast_df["upper"].to_dict()}
        df["พยากรณ์"] = df["ปี_num"].map(fc_mean)
        df["ขอบล่าง"] = df["ปี_num"].map(fc_lower)
        df["ขอบบน"] = df["ปี_num"].map(fc_upper)

        n = len(df)
        step = max(1, round(n / 12))
        tick_vals = df["ปี"].iloc[::step].tolist()
        if df["ปี"].iloc[-1] not in tick_vals:
            tick_vals.append(df["ปี"].iloc[-1])

        x_enc = alt.X(
            "ปี:O", sort=year_order, title=None,
            axis=alt.Axis(values=tick_vals, labelAngle=0, grid=False,
                           domain=False, tickColor="#E9ECF1",
                           labelColor="#5B6B7C", labelFontSize=11, labelPadding=6),
        )
        y_axis = alt.Axis(grid=True, gridColor="#EEF1F5", gridDash=[3, 3],
                           domain=False, tickColor="#E9ECF1",
                           labelColor="#5B6B7C", labelFontSize=11)

        base = alt.Chart(df)

        ci_band = base.mark_area(opacity=0.15, color=forecast_color).encode(
            x=x_enc, y=alt.Y("ขอบล่าง:Q", title=None, axis=y_axis), y2="ขอบบน:Q",
        )
        hist_area = base.mark_area(
            interpolate="monotone", line=False,
            color=alt.Gradient(
                gradient="linear",
                stops=[alt.GradientStop(color=color, offset=0),
                       alt.GradientStop(color="#FFFFFF", offset=1)],
                x1=1, x2=1, y1=1, y2=0,
            ),
            opacity=0.35,
        ).encode(x=x_enc, y=alt.Y("ข้อมูลจริง:Q", title=None, axis=y_axis))
        hist_line = base.mark_line(
            interpolate="monotone", color=color, strokeWidth=2.6,
            point=alt.OverlayMarkDef(filled=True, size=30, color=color, stroke="#FFFFFF", strokeWidth=1.6),
        ).encode(
            x=x_enc, y=alt.Y("ข้อมูลจริง:Q"),
            tooltip=[alt.Tooltip("ปี:O", title="ปี"),
                     alt.Tooltip("ข้อมูลจริง:Q", title="ค่าจริง", format=".4f")],
        )
        fc_line = base.mark_line(
            interpolate="monotone", color=forecast_color, strokeWidth=2.6, strokeDash=[6, 4],
            point=alt.OverlayMarkDef(filled=True, size=30, color=forecast_color, stroke="#FFFFFF", strokeWidth=1.6),
        ).encode(
            x=x_enc, y=alt.Y("พยากรณ์:Q"),
            tooltip=[alt.Tooltip("ปี:O", title="ปี"),
                     alt.Tooltip("พยากรณ์:Q", title="ค่าพยากรณ์", format=".4f")],
        )
        fc_points = base.mark_point(color=forecast_color, filled=True, size=45).transform_filter(
            alt.datum["ปี_num"] > last_year
        ).encode(x=x_enc, y=alt.Y("พยากรณ์:Q"))

        chart = (
            (ci_band + hist_area + hist_line + fc_line + fc_points)
            .properties(height=height, padding={"left": 4, "right": 10, "top": 8, "bottom": 4})
            .configure_view(strokeWidth=0)
            .configure_axis(labelFont=FONT_FAMILY, titleFont=FONT_FAMILY)
        )
        st.altair_chart(chart, use_container_width=True)
        # flex-wrap:wrap กัน legend ตกขอบขวาเวลาหน้าจอแคบ (แทนที่จะโดนตัดหาย
        # ก็ให้มันขึ้นบรรทัดใหม่แทน), row-gap เผื่อกรณีตัดบรรทัด
        st.markdown(
            f'<div style="display:flex;flex-wrap:wrap;column-gap:18px;row-gap:6px;'
            f'font-size:0.82rem;color:var(--brand-navy-soft);margin-top:-6px;">'
            f'<span style="white-space:nowrap;"><span style="display:inline-block;width:10px;height:10px;'
            f'border-radius:50%;background:{color};margin-right:5px;"></span>ข้อมูลจริง</span>'
            f'<span style="white-space:nowrap;"><span style="display:inline-block;width:10px;height:10px;'
            f'border-radius:50%;background:{forecast_color};margin-right:5px;"></span>พยากรณ์ (ARIMA)</span>'
            f'<span style="white-space:nowrap;"><span style="display:inline-block;width:10px;height:10px;'
            f'border-radius:2px;background:{forecast_color};opacity:0.3;margin-right:5px;"></span>'
            f'ช่วงความเชื่อมั่น 95%</span>'
            f'</div>',
            unsafe_allow_html=True,
        )

    if not result_ready:
        st.info("คลิกเพื่อดึงข้อมูลอัตโนมัติจากแถบด้านซ้ายก่อนเพื่อดูกราฟแนวโน้มในหน้านี้")
    else:
        # ================= กราฟภาพรวม: แนวโน้มดัชนี TFP ย้อนหลัง + พยากรณ์ (ARIMA) =================
        st.markdown(
            f'<div class="section-card"><div class="section-title">'
            f'<div class="section-num">{icon("trend-up", 20, 2)}</div>'
            f'<div class="section-title-text"><h3>แนวโน้มดัชนีผลิตภาพการผลิตรวม (TFP) ย้อนหลัง พร้อมพยากรณ์ล่วงหน้า (ARIMA)</h3>'
            f'</div></div>',
            unsafe_allow_html=True,
        )
        tfp_series = model_df[DEP_VAR].dropna().sort_index()
        if tfp_series.empty:
            st.info("ไม่พบข้อมูล TFP ในชุดข้อมูลที่ดึงมา")
        else:
            MIN_POINTS_FOR_ARIMA = 8  # จำนวนปีขั้นต่ำที่พอจะ fit ARIMA ได้อย่างมีความหมาย
            if len(tfp_series) >= MIN_POINTS_FOR_ARIMA:
                fc_col1, fc_col2 = st.columns([3, 1])
                with fc_col1:
                    horizon = st.slider(
                        "จำนวนปีที่ต้องการพยากรณ์ล่วงหน้า", min_value=1, max_value=30,
                        value=5, step=1, key="tfp_forecast_horizon",
                        help="เลือกได้ตั้งแต่ 1 ปีจนถึง 30 ปี ยิ่งพยากรณ์ไกลจากข้อมูลจริง "
                             "ยิ่งมีความไม่แน่นอนสูงขึ้น (ช่วงความเชื่อมั่นจะกว้างขึ้นตามไปด้วย)",
                    )
                with fc_col2:
                    st.markdown("<div style='height:1.9rem;'></div>", unsafe_allow_html=True)
                    st.caption(f"≈ {horizon} ปีข้างหน้า")

                with st.spinner("กำลังหาโมเดล ARIMA ที่เหมาะสมและพยากรณ์..."):
                    forecast_df, arima_order = _auto_arima_forecast(tfp_series, horizon)

                _nice_line_chart_with_forecast(
                    tfp_series, forecast_df, color="#F2811D", forecast_color="#2F6FED", height=340,
                )
                p, d, q = arima_order
                last_fc_year = forecast_df.index.max()
                # ใช้ st.markdown + word-wrap แทน st.caption บรรทัดเดียวยาว ๆ
                # กันข้อความตกขอบขวาในหน้าจอแคบ
                st.markdown(
                    f'<div style="font-size:0.82rem;color:var(--brand-navy-soft);'
                    f'line-height:1.6;overflow-wrap:break-word;">'
                    f'ข้อมูลจริง {len(tfp_series)} ปี (ปี {tfp_series.index.min()}–{tfp_series.index.max()}) '
                    f'| ค่าล่าสุด = {tfp_series.iloc[-1]:.4f} '
                    f'| พยากรณ์ด้วย ARIMA({p},{d},{q}) ถึงปี {last_fc_year} '
                    f'(เลือก order ด้วยค่า AIC ต่ำสุดจากการลอง grid search อัตโนมัติ)'
                    f'</div>',
                    unsafe_allow_html=True,
                )

                with st.expander("📋 ดูตัวเลขพยากรณ์รายปี"):
                    fc_display = forecast_df.rename(
                        columns={"mean": "ค่าพยากรณ์", "lower": "ขอบล่าง 95%", "upper": "ขอบบน 95%"}
                    ).round(4)
                    fc_display.index.name = "ปี"
                    # ตาราง HTML ธีมครีม-ส้ม (คลาส tfp-table-cream) แทน st.dataframe
                    # เดิม เพื่อให้ดีไซน์เข้ากับโทนสีส้ม/ครีมของกราฟพยากรณ์ในส่วนนี้
                    fc_reset = fc_display.reset_index()
                    fc_header_html = "".join(f"<th>{c}</th>" for c in fc_reset.columns)
                    fc_rows_html = "".join(
                        "<tr>" + "".join(
                            f"<td>{int(v) if col == 'ปี' else f'{v:,.4f}'}</td>"
                            for col, v in zip(fc_reset.columns, row)
                        ) + "</tr>"
                        for row in fc_reset.values.tolist()
                    )
                    st.markdown(
                        f'<div style="overflow-x:auto;"><table class="tfp-table-cream"><thead><tr>'
                        f'{fc_header_html}</tr></thead><tbody>{fc_rows_html}</tbody></table></div>',
                        unsafe_allow_html=True,
                    )
                    fc_csv = fc_display.to_csv().encode("utf-8-sig")
                    st.download_button(
                        "ดาวน์โหลดตัวเลขพยากรณ์เป็น CSV",
                        data=fc_csv,
                        file_name="TFP_forecast_ARIMA.csv",
                        mime="text/csv",
                    )
            else:
                _nice_line_chart(tfp_series, color="#F2811D", height=340)
                st.caption(
                    f"ข้อมูล {len(tfp_series)} ปี (ปี {tfp_series.index.min()}–{tfp_series.index.max()}) "
                    f"| ค่าล่าสุด = {tfp_series.iloc[-1]:.4f}"
                )
                st.info(
                    f"ข้อมูลมีเพียง {len(tfp_series)} ปี ยังไม่พอสำหรับพยากรณ์ด้วย ARIMA "
                    f"อย่างน่าเชื่อถือ (ต้องการอย่างน้อย {MIN_POINTS_FOR_ARIMA} ปี)"
                )
        st.markdown('</div>', unsafe_allow_html=True)

        # ================= กราฟรายตัวแปร: แยกกล่องระยะยาว / ระยะสั้น =================
        # แยกรายชื่อตัวแปรอิสระเป็น 2 ชุดตามสมการที่ตัวแปรนั้นอยู่ แทนที่จะรวมเป็น
        # dropdown เดียว — ตัวแปรที่อยู่ในทั้งสองสมการจะไปโผล่ทั้งสองกล่อง (ถูกต้อง
        # เพราะมันมีทั้งผลระยะยาวและระยะสั้นจริง ๆ)
        available_vars_lr = [v for v in active_lr_vars if v != "const" and v in model_df.columns]
        available_vars_sr = [v for v in active_sr_bases if v in model_df.columns]

        def _var_trend_box(title_th: str, options: list, widget_key: str,
                            icon_name: str = "trend-up", accent: str = "#D8A867"):
            """วาดกล่อง selectbox + กราฟเส้นแนวโน้มของตัวแปร 1 ชุด (ยาว หรือ สั้น)
            คืนค่าตัวแปรที่ผู้ใช้เลือกอยู่ในกล่องนี้ (หรือ None ถ้าไม่มีตัวแปรให้เลือก)"""
            # หัวข้อ: แคปซูลขอบสีส้มทั้งสองกล่อง (เดิมแยกส้ม/ฟ้า) + วงกลมไอคอนตันสีส้ม
            # + ตัวหนังสือสีดำ จัดกลาง ตัวใหญ่ขึ้น พร้อมเส้นคาดสีส้มแอบโผล่ใต้แคปซูล
            st.markdown(
                f"""
                <div style="display:flex;justify-content:center;margin-bottom:2px;">
                    <div style="position:relative;">
                        <div style="position:absolute;left:50%;bottom:-5px;transform:translateX(-50%);
                                    width:70%;height:9px;border-radius:6px;background:{accent};
                                    z-index:1;"></div>
                        <div style="position:relative;z-index:2;display:inline-flex;align-items:center;
                                    gap:12px;background:#FFFFFF;border:2px solid {accent};
                                    padding:10px 26px 10px 10px;border-radius:999px;
                                    box-shadow:0 4px 14px rgba(22,50,74,0.08);">
                            <span style="display:flex;align-items:center;justify-content:center;
                                         width:30px;height:30px;border-radius:50%;background:{accent};
                                         color:#FFFFFF;flex-shrink:0;">
                                {icon(icon_name, 16, 2)}
                            </span>
                            <span style="font-weight:700;font-size:0.98rem;color:#000000;
                                         letter-spacing:-0.01em;white-space:nowrap;">{title_th}</span>
                        </div>
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )
            if not options:
                st.info("ไม่พบข้อมูลตัวแปรอิสระในชุดนี้")
                return None
            # กรอบพื้นหลังสีขาวครอบกล่องทั้งก้อน + ช่อง selectbox ตรงกลางสีครีม
            # (สไตล์เดียวกับกล่อง "สมมติตัวแปรเปลี่ยนแปลง" ด้านล่างของหน้านี้)
            st.markdown(
                f"""
                <style>
                .st-key-vartrend_box_{widget_key} {{
                    background: #FFFFFF !important;
                    padding: 22px 22px !important;
                    margin-top: -14px !important;
                }}
                .st-key-vartrend_box_{widget_key} div[data-baseweb="select"] > div {{
                    background: var(--bg-page) !important;
                    border-color: #E7DFCF !important;
                }}
                </style>
                """,
                unsafe_allow_html=True,
            )
            with st.container(border=True, key=f"vartrend_box_{widget_key}"):
                picked = st.selectbox(
                    "เลือกตัวแปรอิสระที่ต้องการดูกราฟ",
                    options=options,
                    format_func=var_label_with_abbr,
                    key=widget_key,
                )
            series = model_df[picked].dropna().sort_index()
            if series.empty:
                st.info("ไม่พบข้อมูลของตัวแปรนี้")
            else:
                _nice_line_chart(series, color="#2F6FED", height=280)
                st.markdown(
                    f'<div style="text-align:center;color:var(--brand-navy-soft);'
                    f'font-size:0.85rem;margin-top:2px;">'
                    f'{var_label_with_abbr(picked)} — ข้อมูล {len(series)} ปี</div>',
                    unsafe_allow_html=True,
                )
            return picked

        st.markdown(
            f'<div class="section-card"><div class="section-title">'
            f'<div class="section-num">{icon("search", 19, 1.8)}</div>'
            f'<div class="section-title-text"><h3>กราฟแนวโน้มตัวแปรอิสระรายตัว</h3>'
            f'</div></div>',
            unsafe_allow_html=True,
        )
        col_lr, col_sr = st.columns(2, gap="large")
        with col_lr:
            chosen_var_lr = _var_trend_box(
                "ตัวแปรอิสระที่ส่งผลระยะยาว", available_vars_lr, "dashboard_var_select_lr",
                icon_name="trend-up", accent="#D8A867",
            )
        with col_sr:
            chosen_var_sr = _var_trend_box(
                "ตัวแปรอิสระที่ส่งผลระยะสั้น", available_vars_sr, "dashboard_var_select_sr",
                icon_name="trend-down", accent="#D8A867",
            )
        st.markdown('</div>', unsafe_allow_html=True)

        # การ์ดผลกระทบด้านล่าง (คำนวณจากค่าความยืดหยุ่น/สัมประสิทธิ์ระยะยาว) ยังคง
        # อ้างอิงตัวแปรที่เลือกในกล่อง "ระยะยาว" เป็นหลัก เพราะสูตร % ผลกระทบสร้างจาก
        # สัมประสิทธิ์ log-log ของสมการระยะยาว (ตัวแปรระยะสั้นจะโชว์เป็นข้อมูลเสริม
        # ในหมายเหตุท้ายการ์ดอยู่แล้วถ้ามีค่าสัมประสิทธิ์ระยะสั้นของตัวแปรเดียวกัน)
        chosen_var = chosen_var_lr
        available_vars = available_vars_lr

        # ================= การ์ดผลกระทบของตัวแปรที่เลือกต่อ TFP (% จากค่าความยืดหยุ่น) =================
        if available_vars:
            st.markdown(
                f'<div class="section-card"><div class="section-title">'
                f'<div class="section-num">{icon("bulb", 20, 1.6)}</div>'
                f'<div class="section-title-text"><h3>ตัวแปรนี้ส่งผลต่อ TFP มากน้อยแค่ไหน</h3>'
                f'</div></div>',
                unsafe_allow_html=True,
            )
            lr_info = lr_raw_map.get(chosen_var)
            sr_info = sr_raw_map.get(chosen_var)
            full_name = _var_full_name(chosen_var)
            is_log = chosen_var.startswith("ln_")

            if lr_info is None:
                st.info(
                    f"'{var_label_with_abbr(chosen_var)}' ไม่ได้อยู่ในสมการระยะยาวชุดปัจจุบัน "
                    f"จึงยังไม่มีค่าสัมประสิทธิ์ให้คำนวณผลกระทบ"
                )
            else:
                coef = lr_info["coef"]
                p_val = lr_info["p"]
                star = _significance_stars(p_val)
                sig_text = {
                    "***": "มีนัยสำคัญทางสถิติสูงมาก (ความเชื่อมั่น 99%)",
                    "**": "มีนัยสำคัญทางสถิติ (ความเชื่อมั่น 95%)",
                    "*": "มีนัยสำคัญทางสถิติเล็กน้อย (ความเชื่อมั่น 90%)",
                    "": "ยังไม่มีนัยสำคัญทางสถิติในระดับที่ยอมรับได้ทั่วไป — ควรตีความตัวเลขด้วยความระมัดระวัง",
                }[star]

                # กรอบพื้นหลังสีขาวครอบกล่อง input ทั้งก้อน + ช่องตัวเลขตรงกลางสีครีม
                # ให้ตัดกับกรอบขาวรอบนอก + ปุ่ม +/- เป็นสีเขียว/แดงค้างไว้ตลอด (ไม่ใช่แค่
                # ตอน hover) ให้ผู้ใช้เห็นชัดว่าปุ่มไหนเพิ่ม/ปุ่มไหนลดค่าโดยไม่ต้องอ่านสัญลักษณ์
                st.markdown(
                    """
                    <style>
                    .st-key-shock_input_box {
                        background: #FFFFFF !important;
                        padding: 22px 22px !important;
                    }
                    div[data-testid="stNumberInput"] div[data-baseweb="input"] {
                        background: var(--bg-page) !important;
                        border: 1.5px solid #E7DFCF !important;
                        border-radius: 8px !important;
                        box-shadow: none !important;
                    }
                    div[data-testid="stNumberInput"] div[data-baseweb="input"]:focus-within {
                        border: 1.5px solid var(--brand-orange) !important;
                        box-shadow: none !important;
                    }
                    div[data-testid="stNumberInput"] input {
                        background: transparent !important;
                        color: var(--brand-navy) !important;
                        font-weight: 700 !important;
                    }
                    /* ปุ่ม + สีเขียว / ปุ่ม − สีแดง ค้างไว้ตลอดเวลา (ไม่ใช่แค่ตอน hover) */
                    button[data-testid="stNumberInputStepUp"],
                    button[aria-label="Increment"] {
                        background: #16A34A !important; border-color: #16A34A !important;
                    }
                    button[data-testid="stNumberInputStepUp"]:hover,
                    button[aria-label="Increment"]:hover {
                        background: #128A3B !important; border-color: #128A3B !important;
                    }
                    button[data-testid="stNumberInputStepUp"] svg,
                    button[aria-label="Increment"] svg { color: #FFFFFF !important; fill: #FFFFFF !important; }
                    button[data-testid="stNumberInputStepDown"],
                    button[aria-label="Decrement"] {
                        background: #EF4444 !important; border-color: #EF4444 !important;
                    }
                    button[data-testid="stNumberInputStepDown"]:hover,
                    button[aria-label="Decrement"]:hover {
                        background: #D63A3A !important; border-color: #D63A3A !important;
                    }
                    button[data-testid="stNumberInputStepDown"] svg,
                    button[aria-label="Decrement"] svg { color: #FFFFFF !important; fill: #FFFFFF !important; }
                    </style>
                    """,
                    unsafe_allow_html=True,
                )

                shock_col, result_col = st.columns([1, 1.3])
                with shock_col:
                    with st.container(border=True, key="shock_input_box"):
                        if is_log:
                            shock = st.number_input(
                                f"สมมติ {full_name} เปลี่ยนแปลง (%)",
                                min_value=-50.0, max_value=50.0, value=1.0, step=0.5,
                                key=f"impact_shock_{chosen_var}",
                            )
                            pct_effect = coef * shock
                            formula_text = (
                                f'<div>ค่าความยืดหยุ่น (elasticity) จากสมการระยะยาว = {coef:.4f}</div>'
                                f'<div style="margin-top:6px;">→ TFP เปลี่ยนแปลง ≈ {coef:.4f} × {shock:g}%</div>'
                            )
                        else:
                            shock = st.number_input(
                                f"สมมติ {full_name} เปลี่ยนแปลง",
                                value=1.0, step=0.5,
                                key=f"impact_shock_{chosen_var}",
                            )
                            pct_effect = (math.exp(coef * shock) - 1) * 100
                            formula_text = (
                                f'<div>สัมประสิทธิ์จากสมการระยะยาว = {coef:.4f} '
                                f'(ตัวแปรนี้ไม่ได้อยู่ในรูป log จึงตีความเป็น semi-elasticity)</div>'
                                f'<div style="margin-top:6px;">→ TFP เปลี่ยนแปลง ≈ '
                                f'(e^({coef:.4f}×{shock:g}) − 1) × 100%</div>'
                            )
                        st.markdown(
                            '<div style="font-size:0.82rem;color:var(--brand-navy-soft);'
                            'line-height:1.5;margin-top:-4px;">'
                            'ปรับตัวเลขด้านบนเพื่อดูว่าถ้าตัวแปรนี้เปลี่ยนแปลงมากน้อยต่างกัน<br>'
                            'TFP จะเปลี่ยนไปกี่ % (คำนวณจากค่าสัมประสิทธิ์ในสมการระยะยาวปัจจุบัน)'
                            '</div>',
                            unsafe_allow_html=True,
                        )

                with result_col:
                    arrow = "↑" if pct_effect >= 0 else "↓"
                    color = "var(--green)" if pct_effect >= 0 else "var(--red)"

                    # ป้ายนัยสำคัญ: สี/ข้อความ ตามระดับดาว (***/**/*/ไม่มี) ให้ดูเป็น badge
                    # เดียวกันแทนที่จะเป็นประโยคลอย ๆ แยกก้อนเหมือนก่อนหน้านี้
                    badge_map = {
                        "***": ("var(--green)", "#E9F9EE", "นัยสำคัญสูงมาก (ความเชื่อมั่น 99%)"),
                        "**": ("var(--green)", "#E9F9EE", "มีนัยสำคัญ (ความเชื่อมั่น 95%)"),
                        "*": ("#B45309", "#FEF3E2", "นัยสำคัญเล็กน้อย (ความเชื่อมั่น 90%)"),
                        "": ("var(--red)", "#FDEDED", "ยังไม่มีนัยสำคัญทางสถิติในระดับที่ยอมรับได้ทั่วไป"),
                    }
                    badge_color, badge_bg, badge_text = badge_map[star]
                    p_text = f"p-value = {p_val:.4f}" if p_val is not None else "ไม่พบค่า p-value"

                    # กรอบขาวเข้ม (border ชัดขึ้น + เงาบาง ๆ) แยกกล่องผลลัพธ์ให้เด่นออกจาก
                    # พื้นหลังหน้าเว็บ แทนกล่องเทาอ่อนแบบเดิมที่กลืนไปกับพื้น
                    st.markdown(
                        f'<div style="background:#FFFFFF;border:1.5px solid #D6DCE5;'
                        f'border-radius:12px;padding:18px 20px;'
                        f'box-shadow:0 1px 4px rgba(15,23,42,0.06);">'
                        f'<div style="font-size:0.78rem;font-weight:600;color:var(--brand-navy-soft);'
                        f'text-transform:uppercase;letter-spacing:.02em;">'
                        f'ผลกระทบต่อดัชนี TFP (สมการระยะยาว)</div>'
                        f'<div style="display:flex;align-items:baseline;gap:8px;margin-top:6px;">'
                        f'<span style="font-size:1.3rem;color:{color};">{arrow}</span>'
                        f'<span style="font-size:2.1rem;font-weight:800;color:{color};line-height:1;">'
                        f'{pct_effect:+.2f}%</span>'
                        f'</div>'
                        f'<div style="font-size:0.78rem;color:var(--brand-navy-soft);margin-top:10px;'
                        f'line-height:1.6;">{formula_text}</div>'
                        f'<div style="height:1px;background:var(--card-border);margin:14px 0 12px;"></div>'
                        f'<div style="display:flex;align-items:center;gap:10px;flex-wrap:wrap;'
                        f'row-gap:8px;">'
                        f'<span style="background:{badge_bg};color:{badge_color};font-size:0.74rem;'
                        f'font-weight:700;padding:5px 12px;border-radius:999px;line-height:1.4;'
                        f'display:inline-block;">'
                        f'{badge_text}</span>'
                        f'<span style="font-size:0.74rem;color:var(--brand-navy-soft);white-space:nowrap;">'
                        f'{p_text}</span>'
                        f'</div>'
                        f'</div>',
                        unsafe_allow_html=True,
                    )

                # แถวหมายเหตุด้านล่าง (short-run + ceteris paribus) รวมเป็นกล่องเดียว
                # กั้นด้วยเส้นประบาง ๆ จากส่วนบน แทนที่จะเป็น caption ลอย ๆ หลายบรรทัด
                footer_rows = []
                if sr_info is not None:
                    footer_rows.append(
                        '<div style="display:flex;gap:10px;">'
                        f'<span style="flex:none;color:var(--brand-orange-dark);">{icon("clock", 14, 1.6)}</span>'
                        '<span><b style="color:var(--brand-navy);">ผลกระทบระยะสั้น (short-run):</b> '
                        f'สัมประสิทธิ์ = {sr_info["coef"]:.4f} — สะท้อนผลของการเปลี่ยนแปลงตัวแปรนี้ในปีนั้น ๆ '
                        'ต่อการเปลี่ยนแปลง TFP ในปีเดียวกัน (คนละความหมายกับผลระยะยาวด้านบน)</span>'
                        '</div>'
                    )
                footer_rows.append(
                    '<div style="display:flex;gap:10px;">'
                    f'<span style="flex:none;color:var(--brand-orange-dark);">{icon("info", 14, 1.6)}</span>'
                    '<span><b style="color:var(--brand-navy);">หมายเหตุ:</b> ตัวเลข % นี้เป็นผลกระทบจากตัวแปรนี้ '
                    '“ตัวเดียว” โดยสมมติให้ตัวแปรอื่นคงที่ (ceteris paribus) ไม่ใช่การพยากรณ์ TFP จริง '
                    'ที่มีหลายปัจจัยเปลี่ยนแปลงพร้อมกัน</span>'
                    '</div>'
                )
                st.markdown(
                    '<div style="margin-top:16px;padding-top:14px;border-top:1px dashed var(--card-border);'
                    'display:flex;flex-direction:column;gap:10px;font-size:0.8rem;'
                    'color:var(--brand-navy-soft);line-height:1.55;">'
                    + "".join(footer_rows) +
                    '</div>',
                    unsafe_allow_html=True,
                )
            st.markdown('</div>', unsafe_allow_html=True)

        # ================= กราฟสัดส่วนอิทธิพลเทียบกันทุกตัวแปร (Standardized Coefficient) =================
        st.markdown(
            f'<div class="section-card"><div class="section-title">'
            f'<div class="section-num">{icon("bars", 19, 1.8)}</div>'
            f'<div class="section-title-text"><h3>สัดส่วนอิทธิพลเทียบกันทุกตัวแปร (%)</h3>'
            f'</div></div>',
            unsafe_allow_html=True,
        )
        influence_vars = [
            v for v in active_lr_vars
            if v != "const" and v in lr_raw_map and v in model_df.columns
        ]
        if len(influence_vars) < 2:
            st.info("ต้องมีตัวแปรอิสระอย่างน้อย 2 ตัวในสมการระยะยาว จึงจะเทียบสัดส่วนอิทธิพลกันได้")
        else:
            sample_df = model_df[[dep_ln] + influence_vars].dropna()
            y_std = sample_df[dep_ln].std() if not sample_df.empty else None
            if sample_df.empty or len(sample_df) < 3 or not y_std or pd.isna(y_std):
                st.info("ข้อมูลไม่พอสำหรับคำนวณสัดส่วนอิทธิพล (ต้องการอย่างน้อย 3 ปีที่มีข้อมูลครบทุกตัวแปร)")
            else:
                rows = []
                for v in influence_vars:
                    x_std = sample_df[v].std()
                    if not x_std or pd.isna(x_std) or x_std == 0:
                        continue
                    std_beta = lr_raw_map[v]["coef"] * (x_std / y_std)
                    rows.append({"code": v, "label": _var_full_name(v), "std_beta": std_beta})

                if not rows:
                    st.info("ไม่สามารถคำนวณสัดส่วนอิทธิพลได้ (ส่วนเบี่ยงเบนมาตรฐานของตัวแปรบางตัวเป็น 0)")
                else:
                    infl_df = pd.DataFrame(rows)
                    infl_df["abs_beta"] = infl_df["std_beta"].abs()
                    total_abs = infl_df["abs_beta"].sum()
                    infl_df["สัดส่วน (%)"] = infl_df["abs_beta"] / total_abs * 100
                    infl_df["ทิศทาง"] = infl_df["std_beta"].apply(
                        lambda x: "หนุนเสริม TFP (+)" if x >= 0 else "ฉุดรั้ง TFP (−)"
                    )
                    infl_df = infl_df.sort_values("สัดส่วน (%)", ascending=False).reset_index(drop=True)

                    dir_scale = alt.Scale(
                        domain=["หนุนเสริม TFP (+)", "ฉุดรั้ง TFP (−)"], range=["#16A34A", "#EF4444"],
                    )
                    y_axis_labels = alt.Axis(domain=False, tickColor="#E9ECF1",
                                              labelColor="#16324A", labelFontSize=12, labelLimit=340)
                    x_axis = alt.Axis(grid=True, gridColor="#EEF1F5", gridDash=[3, 3],
                                       domain=False, tickColor="#E9ECF1",
                                       labelColor="#5B6B7C", labelFontSize=11)

                    # ขยาย domain ของแกน x ให้กว้างกว่าค่ามากที่สุดเล็กน้อย (~18%) กันตัวเลข
                    # ท้ายแท่ง (เช่น 42.0) ที่วางต่อจากแท่งยาวสุดโดนตัดขอบขวาของกราฟ
                    max_pct = float(infl_df["สัดส่วน (%)"].max())
                    pct_scale = alt.Scale(domain=[0, max_pct * 1.18], nice=False)

                    bars = alt.Chart(infl_df).mark_bar(
                        cornerRadiusTopRight=6, cornerRadiusBottomRight=6, height=18,
                    ).encode(
                        x=alt.X("สัดส่วน (%):Q", title=None, axis=x_axis, scale=pct_scale),
                        y=alt.Y("label:N", sort="-x", title=None, axis=y_axis_labels),
                        color=alt.Color("ทิศทาง:N", scale=dir_scale,
                                         legend=alt.Legend(title=None, orient="bottom")),
                        tooltip=[
                            alt.Tooltip("label:N", title="ตัวแปร"),
                            alt.Tooltip("สัดส่วน (%):Q", title="สัดส่วนอิทธิพล", format=".1f"),
                            alt.Tooltip("std_beta:Q", title="Standardized coefficient", format=".4f"),
                            alt.Tooltip("ทิศทาง:N", title="ทิศทาง"),
                        ],
                    )
                    text = alt.Chart(infl_df).mark_text(align="left", dx=5, color="#5B6B7C", fontSize=11).encode(
                        x=alt.X("สัดส่วน (%):Q", scale=pct_scale),
                        y=alt.Y("label:N", sort="-x"),
                        text=alt.Text("สัดส่วน (%):Q", format=".1f"),
                    )
                    chart = (
                        (bars + text)
                        .properties(
                            height=max(220, 42 * len(infl_df)),
                            padding={"left": 15, "right": 15, "top": 5, "bottom": 5},
                            autosize=alt.AutoSizeParams(type="fit", contains="padding"),
                        )
                        .configure_view(strokeWidth=0)
                        .configure_axis(labelFont=FONT_FAMILY, titleFont=FONT_FAMILY)
                        .configure_legend(labelFont=FONT_FAMILY, labelFontSize=12, symbolType="circle")
                    )
                    st.altair_chart(chart, use_container_width=True)
                    st.caption(
                        "คำนวณจาก standardized coefficient (สัมประสิทธิ์ × ส่วนเบี่ยงเบนมาตรฐานของตัวแปรนั้น "
                        "÷ ส่วนเบี่ยงเบนมาตรฐานของ TFP) แล้วนำค่าสัมบูรณ์มาคิดเป็นสัดส่วน % เทียบกันทุกตัวแปร "
                        "ในสมการระยะยาวชุดปัจจุบัน (รวมกันได้ 100%) — เป็นการเทียบ 'น้ำหนักอิทธิพล' ไม่ใช่ "
                        "หน่วยดิบของตัวแปร จึงเทียบข้ามตัวแปรที่มีหน่วยต่างกันได้"
                    )

                    with st.expander("📋 ดูตารางตัวเลข"):
                        table_display = infl_df[["label", "std_beta", "สัดส่วน (%)", "ทิศทาง"]].rename(
                            columns={"label": "ตัวแปร", "std_beta": "Standardized coefficient"}
                        ).round({"Standardized coefficient": 4, "สัดส่วน (%)": 2})
                        # ใช้ตาราง HTML แบบเดียวกับตาราง Diagnostics (คลาส tfp-table)
                        # แทน st.dataframe เพราะ st.dataframe จัดตำแหน่งตัวอักษรราย
                        # คอลัมน์เองไม่ได้ ส่วน tfp-table กำหนด text-align: center
                        # ให้ทุกคอลัมน์ไว้แล้วในสไตล์ชีตด้านบน (ดูตรงคอมเมนต์
                        # "ตาราง HTML สำหรับ Diagnostics")
                        infl_header_html = "".join(f"<th>{c}</th>" for c in table_display.columns)
                        infl_rows_html = "".join(
                            "<tr>" + "".join(f"<td>{v}</td>" for v in row) + "</tr>"
                            for row in table_display.values.tolist()
                        )
                        st.markdown(
                            f'<div style="overflow-x:auto;"><table class="tfp-table"><thead><tr>'
                            f'{infl_header_html}</tr></thead><tbody>{infl_rows_html}</tbody></table></div>',
                            unsafe_allow_html=True,
                        )
                        st.download_button(
                            "ดาวน์โหลดตาราง (.csv)",
                            data=table_display.to_csv(index=False).encode("utf-8-sig"),
                            file_name="TFP_influence_share.csv",
                            mime="text/csv",
                            key="dl_influence_share",
                        )
        st.markdown('</div>', unsafe_allow_html=True)
# ------------------------------------------------------------------------------
# ท้ายหน้าเว็บ — เดิมมี footer กลางหน้าแสดงโลโก้ สอวช./สวค./มหาวิทยาลัย/ภาควิชา +
# ข้อมูลผู้จัดทำอยู่ตรงนี้ ปัจจุบันถูกลบออกแล้ว เพราะย้ายโลโก้มหาวิทยาลัย/ภาควิชา
# และข้อมูลผู้จัดทำทั้งหมดไปแสดงเล็ก ๆ ท้ายแถบเมนูด้านซ้ายแทน (ดูตัวแปร
# _corner_badge_html ที่ประกาศไว้ต้นไฟล์แถวเดียวกับที่กำหนด logo3_html/logo4_html
# และเรียกใช้จริงใน `with st.sidebar:` ด้านบน)
# ------------------------------------------------------------------------------
